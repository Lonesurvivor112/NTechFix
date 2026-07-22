import pandas as pd
import matplotlib.pyplot as plt
from pathlib import Path
from datetime import timedelta
from collections import Counter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# =====================================================
# CONFIG
# =====================================================

OUTPUT_DIR = Path("output")

REPORTS_DIR = Path(
    r"C:\Users\nfalzetti\eisenhowercenter.com"
    r"\IT - Documents\Scripts\Nick Scripts"
    r"\Python\Zammad py\ZammadReports\Reports"
)
REPORTS_DIR.mkdir(parents=True, exist_ok=True)

CHART_DIR = REPORTS_DIR / "tag_report_charts"
CHART_DIR.mkdir(exist_ok=True)

RAW_FILE = OUTPUT_DIR / "Raw_Tickets.xlsx"

LOOKBACK_DAYS = 30

REPORT_DATE = pd.Timestamp.utcnow().tz_localize(None)
PERIOD_START = REPORT_DATE - timedelta(days=LOOKBACK_DAYS)

PERIOD_LABEL = (
    f"{PERIOD_START.strftime('%b %d, %Y')} "
    f"– {REPORT_DATE.strftime('%b %d, %Y')}"
)

FILE_LABEL = REPORT_DATE.strftime("%Y-%m-%d")

PPTX_FILE = (
    REPORTS_DIR /
    f"EC_IT_Tag_Report_{FILE_LABEL}.pptx"
)

REPORT_TITLE = "IT Service Desk Tag Report"
REPORT_SUBTITLE = "Tickets Categorized by Zammad Tag"

TOP_N_TAGS = 15
TICKETS_PER_TAG_TABLE = 10

STATE_OPEN = 2
STATE_CLOSED = 4
STATE_MERGED = 5

# =====================================================
# STYLE
# =====================================================

NAVY = RGBColor(31, 78, 121)
DARK = RGBColor(38, 38, 38)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(221, 235, 247)
LIGHT_GRAY = RGBColor(242, 242, 242)
GREEN_LIGHT = RGBColor(226, 239, 218)
YELLOW_LIGHT = RGBColor(255, 242, 204)
ORANGE_LIGHT = RGBColor(252, 228, 214)
RED_LIGHT = RGBColor(255, 199, 206)
GRAY_LINE = RGBColor(210, 210, 210)
GRAY_TEXT = RGBColor(120, 120, 120)

FONT = "Aptos"

# =====================================================
# HELPERS
# =====================================================

def parse_tags(tag_field):
    if pd.isna(tag_field):
        return []

    if isinstance(tag_field, list):
        raw = tag_field
    else:
        raw = str(tag_field).split(",")

    return [
        t.strip()
        for t in raw
        if t and str(t).strip()
    ]


def add_header(slide, title):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.62)
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    textbox = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.1),
        Inches(12.6), Inches(0.4)
    )

    frame = textbox.text_frame
    frame.clear()

    p = frame.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_footer(slide):
    textbox = slide.shapes.add_textbox(
        Inches(0.35), Inches(7.16),
        Inches(12.6), Inches(0.22)
    )

    frame = textbox.text_frame
    p = frame.paragraphs[0]
    p.text = (
        f"Eisenhower Center IT | "
        f"Tag Report | {PERIOD_LABEL}"
    )
    p.font.name = FONT
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.RIGHT


def add_text(
    slide, text, x, y, w, h,
    size=16, bold=False, color=DARK, align=None
):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    p = frame.paragraphs[0]
    p.text = str(text)
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color

    if align is not None:
        p.alignment = align

    return box


def add_kpi_card(
    slide, label, value, x, y, w, h,
    fill_color=LIGHT_BLUE
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = GRAY_LINE

    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)

    p1 = frame.paragraphs[0]
    p1.text = str(value)
    p1.font.name = FONT
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = NAVY
    p1.alignment = PP_ALIGN.CENTER

    p2 = frame.add_paragraph()
    p2.text = str(label)
    p2.font.name = FONT
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK
    p2.alignment = PP_ALIGN.CENTER


def add_table(
    slide, dataframe, x, y, w, h,
    font_size=10, max_rows=10
):
    if dataframe.empty:
        add_text(
            slide, "No data available.",
            x, y, w, h, size=12
        )
        return

    df = dataframe.head(max_rows).copy()

    rows = len(df) + 1
    cols = len(df.columns)

    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    table = table_shape.table

    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = FONT
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = WHITE

    for row_idx, (_, row) in enumerate(
        df.iterrows(), start=1
    ):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)

            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = DARK


def create_bar_chart(
    df, label_col, value_col, title, filename,
    top_n=15, horizontal=True
):
    chart_path = CHART_DIR / filename

    if df.empty:
        return None

    if label_col not in df.columns or value_col not in df.columns:
        return None

    data = df.head(top_n).copy()

    plt.figure(figsize=(9, 5.5))

    if horizontal:
        data = data.sort_values(
            value_col, ascending=True
        )
        plt.barh(
            data[label_col].astype(str),
            data[value_col]
        )
        plt.xlabel(value_col)
    else:
        plt.bar(
            data[label_col].astype(str),
            data[value_col]
        )
        plt.ylabel(value_col)
        plt.xticks(rotation=35, ha="right")

    plt.title(title)
    plt.tight_layout()
    plt.savefig(chart_path, dpi=160)
    plt.close()

    return chart_path


def create_grouped_bar_chart(
    df, label_col, cols, title, filename, top_n=15
):
    chart_path = CHART_DIR / filename

    if df.empty:
        return None

    data = df.head(top_n).copy()
    data = data.sort_values(cols[0], ascending=True)

    labels = data[label_col].astype(str)

    fig, ax = plt.subplots(figsize=(10, 5.5))

    y_pos = range(len(labels))
    height = 0.4

    for i, col in enumerate(cols):
        offset = (i - (len(cols) - 1) / 2) * height
        positions = [y + offset for y in y_pos]
        ax.barh(positions, data[col], height=height, label=col)

    ax.set_yticks(list(y_pos))
    ax.set_yticklabels(labels)
    ax.set_xlabel("Tickets")
    ax.set_title(title)
    ax.legend()

    plt.tight_layout()
    plt.savefig(chart_path, dpi=160)
    plt.close()

    return chart_path


def add_chart_image(slide, chart_path, x, y, w, h):
    if chart_path and Path(chart_path).exists():
        slide.shapes.add_picture(
            str(chart_path),
            Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
    else:
        add_text(
            slide,
            "Chart could not be generated.",
            x, y, w, h, size=12
        )


def build_tag_counter(dataframe):
    counter = Counter()

    for tags in dataframe["ParsedTags"]:
        for tag in tags:
            counter[tag] += 1

    return counter


def counter_to_df(counter, tag_col, count_col):
    return pd.DataFrame(
        counter.most_common(),
        columns=[tag_col, count_col]
    )


# =====================================================
# LOAD DATA
# =====================================================

print(f"Reporting Date: {REPORT_DATE.date()}")
print(f"Recent Window: {PERIOD_LABEL}")
print()
print("Loading tickets...")

df = pd.read_excel(RAW_FILE)

print(f"Loaded {len(df):,} tickets")

for col in ["created_at", "updated_at", "close_at"]:
    if col in df.columns:
        df[col] = pd.to_datetime(
            df[col], errors="coerce", utc=True
        )
        df[col] = df[col].dt.tz_localize(None)

if "state_id" in df.columns:
    df["state_id"] = pd.to_numeric(
        df["state_id"], errors="coerce"
    )

    merged_count = len(df[df["state_id"] == STATE_MERGED])
    df = df[df["state_id"] != STATE_MERGED].copy()

    print(f"Excluded {merged_count:,} merged tickets")

if "article_count" in df.columns:
    df["article_count"] = pd.to_numeric(
        df["article_count"], errors="coerce"
    ).fillna(0)

if "time_unit" not in df.columns:
    df["time_unit"] = 0

df["time_unit"] = pd.to_numeric(
    df["time_unit"], errors="coerce"
).fillna(0)

if "tags" not in df.columns:
    print(
        "ERROR: No 'tags' column. Run Enrich-ZammadTags.py first."
    )
    raise SystemExit(1)

df["ParsedTags"] = df["tags"].apply(parse_tags)

# =====================================================
# SEGMENTS
# =====================================================

tagged_df = df[df["ParsedTags"].apply(len) > 0].copy()

recent_df = df[
    (df["created_at"] >= PERIOD_START) &
    (df["created_at"] <= REPORT_DATE)
].copy()

if "state_id" in df.columns:
    open_df = df[df["state_id"] == STATE_OPEN].copy()
else:
    open_df = df[df["close_at"].isna()].copy()

open_df["Days Open"] = (
    (REPORT_DATE - open_df["created_at"])
    .dt.total_seconds()
    .div(86400)
    .fillna(0)
    .astype(int)
    .clip(lower=0)
)

# =====================================================
# TAG METRICS
# =====================================================

print()
print("Building tag counters...")

all_tags_counter = build_tag_counter(df)
recent_tags_counter = build_tag_counter(recent_df)
open_tags_counter = build_tag_counter(open_df)

all_tags_df = counter_to_df(
    all_tags_counter, "Tag", "Total Tickets"
)
recent_tags_df = counter_to_df(
    recent_tags_counter, "Tag", "Last 30 Days"
)
open_tags_df = counter_to_df(
    open_tags_counter, "Tag", "Open Tickets"
)

combined = all_tags_df.merge(
    recent_tags_df, on="Tag", how="left"
).merge(
    open_tags_df, on="Tag", how="left"
).fillna(0)

for col in ["Last 30 Days", "Open Tickets"]:
    combined[col] = combined[col].astype(int)

combined = combined.sort_values(
    "Total Tickets", ascending=False
)

# =====================================================
# TAG TIME SPENT
# =====================================================

print("Building tag time spent...")

tag_time_rows = []

for _, ticket in df.iterrows():
    minutes = ticket.get("time_unit", 0) or 0

    for tag in ticket["ParsedTags"]:
        tag_time_rows.append({
            "Tag": tag,
            "Minutes": minutes
        })

tag_time_raw = pd.DataFrame(tag_time_rows)

if not tag_time_raw.empty:
    tag_time = (
        tag_time_raw
        .groupby("Tag")
        .agg(
            Tickets=("Minutes", "count"),
            TicketsWithTime=(
                "Minutes",
                lambda s: int((s > 0).sum())
            ),
            TotalMinutes=("Minutes", "sum"),
            AvgMinutesPerTicket=("Minutes", "mean")
        )
        .reset_index()
    )

    tag_time["AvgMinutesPerTicket"] = (
        tag_time["AvgMinutesPerTicket"].round(1)
    )
    tag_time["TotalHours"] = (
        (tag_time["TotalMinutes"] / 60).round(1)
    )
    tag_time["TimeCoverage %"] = (
        (
            tag_time["TicketsWithTime"] /
            tag_time["Tickets"] * 100
        ).round(1)
    )

    tag_time = tag_time.sort_values(
        "TotalMinutes", ascending=False
    )
else:
    tag_time = pd.DataFrame(
        columns=[
            "Tag", "Tickets", "TicketsWithTime",
            "TotalMinutes", "AvgMinutesPerTicket",
            "TotalHours", "TimeCoverage %"
        ]
    )

# =====================================================
# EXPORT XLSX
# =====================================================

print("Saving XLSX outputs...")

combined.to_excel(
    REPORTS_DIR / f"Tag_Report_Combined_{FILE_LABEL}.xlsx",
    index=False
)

tag_time.to_excel(
    REPORTS_DIR / f"Tag_Report_TimeSpent_{FILE_LABEL}.xlsx",
    index=False
)

# Per-tag ticket detail
print("Building per-tag ticket lists...")

top_tags_for_detail = list(
    all_tags_df.head(TOP_N_TAGS)["Tag"]
)

per_tag_rows = []

for tag in top_tags_for_detail:
    matching = df[
        df["ParsedTags"].apply(
            lambda tags: tag in tags
        )
    ].copy()

    if matching.empty:
        continue

    if "state_id" in matching.columns:
        matching["Status"] = matching["state_id"].map({
            STATE_OPEN: "Open",
            STATE_CLOSED: "Closed"
        }).fillna("Other")
    else:
        matching["Status"] = matching["close_at"].apply(
            lambda v: "Closed" if pd.notna(v) else "Open"
        )

    matching["Tag"] = tag

    per_tag_rows.append(matching)

if per_tag_rows:
    per_tag_all = pd.concat(per_tag_rows, ignore_index=True)

    export_cols = [
        c for c in [
            "Tag", "number", "title", "Status",
            "created_at", "close_at", "time_unit",
            "state_id"
        ]
        if c in per_tag_all.columns
    ]

    per_tag_all[export_cols].to_excel(
        REPORTS_DIR / f"Tag_Report_TicketDetail_{FILE_LABEL}.xlsx",
        index=False
    )

# =====================================================
# CHARTS
# =====================================================

print("Creating charts...")

chart_all_tags = create_bar_chart(
    all_tags_df,
    "Tag", "Total Tickets",
    "Top Tags — All Time",
    "top_tags_all_time.png",
    top_n=TOP_N_TAGS
)

chart_recent_tags = create_bar_chart(
    recent_tags_df,
    "Tag", "Last 30 Days",
    "Top Tags — Last 30 Days",
    "top_tags_recent.png",
    top_n=TOP_N_TAGS
)

chart_open_tags = create_bar_chart(
    open_tags_df,
    "Tag", "Open Tickets",
    "Top Tags — Currently Open",
    "top_tags_open.png",
    top_n=TOP_N_TAGS
)

chart_time_tags = create_bar_chart(
    tag_time,
    "Tag", "TotalHours",
    "Top Tags by Time Logged (Hours)",
    "top_tags_time.png",
    top_n=TOP_N_TAGS
)

chart_combined = create_grouped_bar_chart(
    combined,
    "Tag",
    ["Total Tickets", "Last 30 Days", "Open Tickets"],
    "Tag Comparison: All-Time vs Recent vs Open",
    "tag_comparison.png",
    top_n=TOP_N_TAGS
)

# =====================================================
# BUILD PRESENTATION
# =====================================================

print("Building PowerPoint...")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

# --- Slide 1: Title ---

slide = prs.slides.add_slide(blank_layout)

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(13.333), Inches(7.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = NAVY
shape.line.fill.background()

add_text(
    slide, REPORT_TITLE,
    0.8, 2.25, 11.8, 0.7,
    size=34, bold=True, color=WHITE
)
add_text(
    slide, REPORT_SUBTITLE,
    0.83, 3.05, 11.8, 0.4,
    size=18, color=WHITE
)
add_text(
    slide, PERIOD_LABEL,
    0.83, 3.65, 11.8, 0.4,
    size=14, color=WHITE
)
add_text(
    slide, "Eisenhower Center IT",
    0.83, 5.95, 11.8, 0.3,
    size=14, color=WHITE
)

# --- Slide 2: Overview ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Tag Report Overview")
add_footer(slide)

total_tickets = len(df)
tagged_tickets = len(tagged_df)
recent_tickets = len(recent_df)
open_tickets = len(open_df)
unique_tags = len(all_tags_counter)

coverage_pct = (
    round(tagged_tickets / total_tickets * 100, 1)
    if total_tickets else 0
)

if not all_tags_df.empty:
    top_all_tag = all_tags_df.iloc[0]["Tag"]
else:
    top_all_tag = "N/A"

if not recent_tags_df.empty:
    top_recent_tag = recent_tags_df.iloc[0]["Tag"]
else:
    top_recent_tag = "N/A"

if not tag_time.empty:
    top_time_tag = tag_time.iloc[0]["Tag"]
    top_time_hours = tag_time.iloc[0]["TotalHours"]
    top_time_str = f"{top_time_tag} ({top_time_hours} hrs)"
else:
    top_time_str = "N/A"

add_kpi_card(
    slide, "Total Tickets", total_tickets,
    0.45, 1.0, 2.0, 1.15
)
add_kpi_card(
    slide, "Tagged Tickets", tagged_tickets,
    2.60, 1.0, 2.0, 1.15,
    fill_color=GREEN_LIGHT
)
add_kpi_card(
    slide, "Tag Coverage %", f"{coverage_pct}%",
    4.75, 1.0, 2.0, 1.15,
    fill_color=YELLOW_LIGHT
)
add_kpi_card(
    slide, "Unique Tags", unique_tags,
    6.90, 1.0, 2.0, 1.15
)
add_kpi_card(
    slide, "Top Tag (All-Time)", top_all_tag,
    9.05, 1.0, 2.0, 1.15,
    fill_color=LIGHT_BLUE
)
add_kpi_card(
    slide, "Top Tag by Time", top_time_tag if not tag_time.empty else "N/A",
    11.20, 1.0, 1.9, 1.15,
    fill_color=ORANGE_LIGHT
)

add_text(
    slide, "How to Read This Report",
    0.55, 2.55, 12.5, 0.4,
    size=18, bold=True, color=NAVY
)

blurbs = [
    "- All-Time — every Zammad ticket, regardless of when it was opened.",
    "- Last 30 Days — tickets created in the current reporting window.",
    "- Currently Open — tickets still in open status right now.",
    "- Time Logged — real minutes captured via Zammad's time_unit field, aggregated per tag.",
    "- Tag Coverage — % of tickets with at least one Zammad tag.",
    "- Older tickets pre-date tagging discipline and land in 'Other' via category fallback."
]

y = 3.05
for line in blurbs:
    add_text(
        slide, line,
        0.75, y, 12.2, 0.45,
        size=13, color=DARK
    )
    y += 0.45

# --- Slide 3: All-Time ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Top Tags — All Time")
add_footer(slide)

add_chart_image(
    slide, chart_all_tags,
    0.55, 1.0, 7.6, 5.9
)

add_table(
    slide, all_tags_df.head(TOP_N_TAGS),
    8.45, 1.05, 4.3, 5.85,
    font_size=10, max_rows=TOP_N_TAGS
)

# --- Slide 4: Recent ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Top Tags — Last 30 Days")
add_footer(slide)

add_chart_image(
    slide, chart_recent_tags,
    0.55, 1.0, 7.6, 5.9
)

add_table(
    slide, recent_tags_df.head(TOP_N_TAGS),
    8.45, 1.05, 4.3, 5.85,
    font_size=10, max_rows=TOP_N_TAGS
)

# --- Slide 5: Open ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Currently Open Tags")
add_footer(slide)

add_chart_image(
    slide, chart_open_tags,
    0.55, 1.0, 7.6, 5.9
)

add_table(
    slide, open_tags_df.head(TOP_N_TAGS),
    8.45, 1.05, 4.3, 5.85,
    font_size=10, max_rows=TOP_N_TAGS
)

# --- Slide 6: Comparison ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Tag Trends — All-Time vs Recent vs Open")
add_footer(slide)

add_chart_image(
    slide, chart_combined,
    0.55, 1.0, 12.2, 5.9
)

# --- Slide 7: Time by Tag ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Time Logged by Tag")
add_footer(slide)

add_chart_image(
    slide, chart_time_tags,
    0.55, 1.0, 7.6, 5.9
)

time_display_cols = [
    c for c in [
        "Tag", "Tickets", "TicketsWithTime",
        "TotalHours", "AvgMinutesPerTicket"
    ]
    if c in tag_time.columns
]

add_table(
    slide, tag_time[time_display_cols].head(TOP_N_TAGS),
    8.35, 1.05, 4.65, 5.85,
    font_size=8, max_rows=TOP_N_TAGS
)

# --- Slides 8+: Per-Tag Detail ---

for tag in top_tags_for_detail:
    matching = df[
        df["ParsedTags"].apply(
            lambda tags: tag in tags
        )
    ].copy()

    if matching.empty:
        continue

    if "state_id" in matching.columns:
        matching["Status"] = matching["state_id"].map({
            STATE_OPEN: "Open",
            STATE_CLOSED: "Closed"
        }).fillna("Other")
    else:
        matching["Status"] = matching["close_at"].apply(
            lambda v: "Closed" if pd.notna(v) else "Open"
        )

    total_tag = len(matching)
    open_tag = (matching["Status"] == "Open").sum()
    closed_tag = (matching["Status"] == "Closed").sum()

    recent_tag = matching[
        (matching["created_at"] >= PERIOD_START) &
        (matching["created_at"] <= REPORT_DATE)
    ]
    recent_count = len(recent_tag)

    total_minutes = int(matching["time_unit"].sum())
    total_hours = round(total_minutes / 60, 1)

    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, f"Tag Detail: {tag}")
    add_footer(slide)

    add_kpi_card(
        slide, "Total Tickets", total_tag,
        0.45, 1.0, 2.4, 1.05
    )
    add_kpi_card(
        slide, "Closed", closed_tag,
        3.05, 1.0, 2.4, 1.05,
        fill_color=GREEN_LIGHT
    )
    add_kpi_card(
        slide, "Currently Open", open_tag,
        5.65, 1.0, 2.4, 1.05,
        fill_color=ORANGE_LIGHT
    )
    add_kpi_card(
        slide, "Last 30 Days", recent_count,
        8.25, 1.0, 2.4, 1.05,
        fill_color=YELLOW_LIGHT
    )
    add_kpi_card(
        slide, "Time Logged (hrs)", total_hours,
        10.85, 1.0, 2.0, 1.05,
        fill_color=RED_LIGHT
    )

    detail_cols = [
        c for c in [
            "number", "title", "Status",
            "created_at", "time_unit"
        ]
        if c in matching.columns
    ]

    display = (
        matching[detail_cols]
        .sort_values(
            by="created_at",
            ascending=False,
            na_position="last"
        )
        .head(TICKETS_PER_TAG_TABLE)
        .copy()
    )

    if "created_at" in display.columns:
        display["created_at"] = (
            display["created_at"]
            .dt.strftime("%Y-%m-%d")
        )
        display.rename(
            columns={"created_at": "Created"},
            inplace=True
        )

    if "time_unit" in display.columns:
        display.rename(
            columns={"time_unit": "Minutes"},
            inplace=True
        )

    add_text(
        slide,
        f"Most recent {TICKETS_PER_TAG_TABLE} tickets tagged '{tag}':",
        0.55, 2.35, 12.2, 0.35,
        size=13, bold=True, color=NAVY
    )

    add_table(
        slide, display,
        0.55, 2.75, 12.2, 4.3,
        font_size=10,
        max_rows=TICKETS_PER_TAG_TABLE
    )

# --- Final Slide: Notes ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Notes on Tag Data")
add_footer(slide)

notes = [
    f"Report generated {REPORT_DATE.strftime('%B %d, %Y')} at Eisenhower Center.",
    f"Tag coverage is currently {coverage_pct}% — older tickets pre-date the tagging system.",
    "Tag Report is complementary to the Executive/Monthly Reports — tags are raw Zammad values, categories are grouped for management view.",
    "Merged tickets (state_id=5) are excluded from all counts.",
    "Time Logged is the sum of Zammad time_unit values, which agents record per ticket. Coverage may vary by tag.",
    "Improving tagging and time-logging discipline in Zammad will directly improve future report accuracy."
]

y = 1.1
for idx, note in enumerate(notes, start=1):
    add_text(
        slide, f"{idx}. {note}",
        0.75, y, 11.9, 0.6,
        size=14, color=DARK
    )
    y += 0.65

# =====================================================
# SAVE
# =====================================================

prs.save(PPTX_FILE)

print()
print("=" * 60)
print("TAG REPORT CREATED")
print("=" * 60)
print(PPTX_FILE)
print()
print("Done.")