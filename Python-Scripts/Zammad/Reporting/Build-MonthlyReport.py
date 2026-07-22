import json
import pandas as pd
import matplotlib.pyplot as plt
from pathlib import Path
from datetime import timedelta
from collections import Counter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# =====================================================
# CONFIG
# =====================================================

OUTPUT_DIR = Path("output")

REPORTS_DIR = Path(
    r"C:\Users\nfalzetti\eisenhowercenter.com"
    r"\IT - Documents\Scripts\Nick Scripts"
    r"\Python\Zammad py\ZammadReports\Reports"
)
REPORTS_DIR.mkdir(parents=True, exist_ok=True)

CHART_DIR = REPORTS_DIR / "monthly_charts"
CHART_DIR.mkdir(exist_ok=True)

RAW_FILE = OUTPUT_DIR / "Raw_Tickets.xlsx"

LOOKBACK_DAYS = 30

REPORT_DATE = pd.Timestamp.utcnow().tz_localize(None)
PERIOD_START = REPORT_DATE - timedelta(days=LOOKBACK_DAYS)

PERIOD_LABEL = (
    f"{PERIOD_START.strftime('%b %d, %Y')} "
    f"– {REPORT_DATE.strftime('%b %d, %Y')}"
)

FILE_LABEL = REPORT_DATE.strftime("%Y-%m")

PPTX_FILE = (
    REPORTS_DIR /
    f"EC_IT_Monthly_Report_{FILE_LABEL}.pptx"
)

REPORT_TITLE = "IT Service Desk Monthly Report"
REPORT_SUBTITLE = f"Reporting Period: {PERIOD_LABEL}"

CATEGORY_FILES = {
    "keyword": "config/categories.json",
    "tag": "config/tag_categories.json"
}

STATE_OPEN = 2
STATE_CLOSED = 4
STATE_MERGED = 5

# =====================================================
# STYLE
# =====================================================

NAVY = RGBColor(31, 78, 121)
DARK = RGBColor(38, 38, 38)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(221, 235, 247)
LIGHT_GRAY = RGBColor(242, 242, 242)
GREEN_LIGHT = RGBColor(226, 239, 218)
YELLOW_LIGHT = RGBColor(255, 242, 204)
ORANGE_LIGHT = RGBColor(252, 228, 214)
RED_LIGHT = RGBColor(255, 199, 206)
GRAY_LINE = RGBColor(210, 210, 210)
GRAY_TEXT = RGBColor(120, 120, 120)

FONT = "Aptos"

# =====================================================
# CATEGORY DEFS
# =====================================================

with open(
    CATEGORY_FILES["keyword"],
    "r", encoding="utf-8"
) as f:
    KEYWORD_CATEGORIES = json.load(f)

with open(
    CATEGORY_FILES["tag"],
    "r", encoding="utf-8"
) as f:
    TAG_CATEGORIES = json.load(f)


def parse_tags(tag_field):
    if pd.isna(tag_field):
        return []

    if isinstance(tag_field, list):
        raw = tag_field
    else:
        raw = str(tag_field).split(",")

    return [
        t.strip()
        for t in raw
        if t and str(t).strip()
    ]


def categorize_by_tags(tag_list):
    tag_set = {
        t.strip().lower()
        for t in tag_list
    }

    for category, tags in TAG_CATEGORIES.items():
        for tag in tags:
            if tag.strip().lower() in tag_set:
                return category

    return None


def categorize_by_keywords(title):
    title = str(title).lower().strip()

    best_category = "Other"
    best_match_len = 0

    for category, keywords in KEYWORD_CATEGORIES.items():
        for keyword in keywords:
            kw = str(keyword).lower().strip()

            if kw and kw in title:
                if len(kw) > best_match_len:
                    best_match_len = len(kw)
                    best_category = category

    return best_category


def categorize_row(row):
    tag_list = parse_tags(row.get("tags"))

    if tag_list:
        tag_category = categorize_by_tags(tag_list)

        if tag_category:
            return tag_category

    return categorize_by_keywords(row.get("title", ""))


# =====================================================
# PPTX HELPERS
# =====================================================

def add_header(slide, title):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.62)
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    textbox = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.1),
        Inches(12.6), Inches(0.4)
    )

    frame = textbox.text_frame
    frame.clear()

    p = frame.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_footer(slide):
    textbox = slide.shapes.add_textbox(
        Inches(0.35), Inches(7.16),
        Inches(12.6), Inches(0.22)
    )

    frame = textbox.text_frame
    p = frame.paragraphs[0]
    p.text = (
        f"Eisenhower Center IT | "
        f"Monthly Report | {PERIOD_LABEL}"
    )
    p.font.name = FONT
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.RIGHT


def add_text(
    slide, text, x, y, w, h,
    size=16, bold=False, color=DARK, align=None
):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    p = frame.paragraphs[0]
    p.text = str(text)
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color

    if align is not None:
        p.alignment = align

    return box


def add_kpi_card(
    slide, label, value, x, y, w, h,
    fill_color=LIGHT_BLUE
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = GRAY_LINE

    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)

    p1 = frame.paragraphs[0]
    p1.text = str(value)
    p1.font.name = FONT
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = NAVY
    p1.alignment = PP_ALIGN.CENTER

    p2 = frame.add_paragraph()
    p2.text = str(label)
    p2.font.name = FONT
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK
    p2.alignment = PP_ALIGN.CENTER


def add_table(
    slide, dataframe, x, y, w, h,
    font_size=10, max_rows=10
):
    if dataframe.empty:
        add_text(
            slide, "No data available.",
            x, y, w, h, size=12
        )
        return

    df = dataframe.head(max_rows).copy()

    rows = len(df) + 1
    cols = len(df.columns)

    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    table = table_shape.table

    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = FONT
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = WHITE

    for row_idx, (_, row) in enumerate(
        df.iterrows(), start=1
    ):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)

            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = DARK


def create_bar_chart(
    df, label_col, value_col, title, filename,
    top_n=10, horizontal=True
):
    chart_path = CHART_DIR / filename

    if df.empty:
        return None

    if label_col not in df.columns or value_col not in df.columns:
        return None

    data = df.head(top_n).copy()

    plt.figure(figsize=(9, 4.8))

    if horizontal:
        data = data.sort_values(
            value_col, ascending=True
        )
        plt.barh(
            data[label_col].astype(str),
            data[value_col]
        )
        plt.xlabel(value_col)
    else:
        plt.bar(
            data[label_col].astype(str),
            data[value_col]
        )
        plt.ylabel(value_col)
        plt.xticks(rotation=35, ha="right")

    plt.title(title)
    plt.tight_layout()
    plt.savefig(chart_path, dpi=160)
    plt.close()

    return chart_path


def create_daily_trend_chart(df, filename):
    chart_path = CHART_DIR / filename

    if df.empty:
        return None

    plt.figure(figsize=(10, 4.5))

    plt.plot(
        df["Date"], df["Created"],
        marker="o", label="Created"
    )
    plt.plot(
        df["Date"], df["Closed"],
        marker="o", label="Closed"
    )

    plt.title("Daily Ticket Activity (Last 30 Days)")
    plt.xlabel("Date")
    plt.ylabel("Tickets")
    plt.xticks(rotation=45, ha="right")
    plt.legend()
    plt.tight_layout()
    plt.savefig(chart_path, dpi=160)
    plt.close()

    return chart_path


def add_chart_image(slide, chart_path, x, y, w, h):
    if chart_path and Path(chart_path).exists():
        slide.shapes.add_picture(
            str(chart_path),
            Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
    else:
        add_text(
            slide,
            "Chart could not be generated.",
            x, y, w, h, size=12
        )


# =====================================================
# LOAD DATA
# =====================================================

print(f"Reporting Period: {PERIOD_LABEL}")
print()
print("Loading tickets...")

df = pd.read_excel(RAW_FILE)

print(f"Loaded {len(df):,} tickets")

for col in ["created_at", "updated_at", "close_at"]:
    if col in df.columns:
        df[col] = pd.to_datetime(
            df[col], errors="coerce", utc=True
        )
        df[col] = df[col].dt.tz_localize(None)

if "state_id" in df.columns:
    df["state_id"] = pd.to_numeric(
        df["state_id"], errors="coerce"
    )

    merged_count = len(df[df["state_id"] == STATE_MERGED])
    df = df[df["state_id"] != STATE_MERGED].copy()

    print(f"Excluded {merged_count:,} merged tickets")

if "article_count" in df.columns:
    df["article_count"] = pd.to_numeric(
        df["article_count"], errors="coerce"
    ).fillna(0)

if "time_unit" not in df.columns:
    df["time_unit"] = 0

df["time_unit"] = pd.to_numeric(
    df["time_unit"], errors="coerce"
).fillna(0)

if "tags" not in df.columns:
    df["tags"] = ""

print("Categorizing tickets...")

df["ParsedTags"] = df["tags"].apply(parse_tags)
df["Category"] = df.apply(categorize_row, axis=1)

# =====================================================
# WINDOW SLICES
# =====================================================

created_this_month = df[
    (df["created_at"] >= PERIOD_START) &
    (df["created_at"] <= REPORT_DATE)
].copy()

closed_this_month = df[
    (df["close_at"] >= PERIOD_START) &
    (df["close_at"] <= REPORT_DATE) &
    (df["state_id"] == STATE_CLOSED)
].copy()

open_now = df[df["state_id"] == STATE_OPEN].copy()

open_now["Days Open"] = (
    (REPORT_DATE - open_now["created_at"])
    .dt.total_seconds()
    .div(86400)
    .fillna(0)
    .astype(int)
    .clip(lower=0)
)

aged_backlog = len(
    open_now[open_now["Days Open"] > 30]
)

# =====================================================
# METRICS
# =====================================================

print("Building monthly metrics...")

top_categories_month = (
    created_this_month["Category"]
    .value_counts()
    .reset_index()
)
top_categories_month.columns = ["Category", "Tickets"]

tag_counter = Counter()

for tags in created_this_month["ParsedTags"]:
    for tag in tags:
        tag_counter[tag] += 1

top_tags_month = pd.DataFrame(
    tag_counter.most_common(),
    columns=["Tag", "Tickets"]
)

open_by_category = (
    open_now["Category"]
    .value_counts()
    .reset_index()
)
open_by_category.columns = ["Category", "Open Tickets"]

# Time-spent this month by category
time_by_category_month = (
    created_this_month
    .groupby("Category")
    .agg(
        Tickets=("id", "count"),
        TotalMinutes=("time_unit", "sum"),
        TicketsWithTime=(
            "time_unit",
            lambda s: int((s > 0).sum())
        )
    )
    .reset_index()
)
time_by_category_month["TotalHours"] = (
    (time_by_category_month["TotalMinutes"] / 60).round(1)
)
time_by_category_month = time_by_category_month.sort_values(
    "TotalMinutes", ascending=False
)

# Time-spent this month by tag
tag_time_rows = []
for _, ticket in created_this_month.iterrows():
    minutes = ticket.get("time_unit", 0) or 0

    for tag in ticket["ParsedTags"]:
        tag_time_rows.append({
            "Tag": tag,
            "Minutes": minutes
        })

tag_time_df = pd.DataFrame(tag_time_rows)

if not tag_time_df.empty:
    time_by_tag_month = (
        tag_time_df
        .groupby("Tag")
        .agg(
            Tickets=("Minutes", "count"),
            TotalMinutes=("Minutes", "sum")
        )
        .reset_index()
    )
    time_by_tag_month["TotalHours"] = (
        (time_by_tag_month["TotalMinutes"] / 60).round(1)
    )
    time_by_tag_month = time_by_tag_month.sort_values(
        "TotalMinutes", ascending=False
    )
else:
    time_by_tag_month = pd.DataFrame(
        columns=["Tag", "Tickets", "TotalMinutes", "TotalHours"]
    )

# Daily trend
daily_created = (
    created_this_month.groupby(
        created_this_month["created_at"].dt.date
    )
    .size()
    .reset_index(name="Created")
    .rename(columns={"created_at": "Date"})
)

daily_closed = (
    closed_this_month.groupby(
        closed_this_month["close_at"].dt.date
    )
    .size()
    .reset_index(name="Closed")
    .rename(columns={"close_at": "Date"})
)

daily_trend = daily_created.merge(
    daily_closed,
    on="Date",
    how="outer"
).fillna(0)

daily_trend["Created"] = daily_trend["Created"].astype(int)
daily_trend["Closed"] = daily_trend["Closed"].astype(int)
daily_trend = daily_trend.sort_values("Date")

# High risk
high_risk_month = (
    open_now[open_now["Days Open"] > 30]
    .sort_values("Days Open", ascending=False)
    .copy()
)

# Aggregate metrics
total_minutes_month = int(
    created_this_month["time_unit"].sum()
)
total_hours_month = round(total_minutes_month / 60, 1)

tickets_with_time_month = int(
    (created_this_month["time_unit"] > 0).sum()
)

if len(created_this_month) > 0:
    time_coverage_month = round(
        tickets_with_time_month /
        len(created_this_month) * 100,
        1
    )
else:
    time_coverage_month = 0

# =====================================================
# EXPORT DATA
# =====================================================

print("Saving monthly XLSX outputs...")

top_categories_month.to_excel(
    REPORTS_DIR / f"Monthly_Top_Categories_{FILE_LABEL}.xlsx",
    index=False
)

top_tags_month.to_excel(
    REPORTS_DIR / f"Monthly_Top_Tags_{FILE_LABEL}.xlsx",
    index=False
)

open_by_category.to_excel(
    REPORTS_DIR / f"Monthly_Open_Categories_{FILE_LABEL}.xlsx",
    index=False
)

daily_trend.to_excel(
    REPORTS_DIR / f"Monthly_Daily_Trend_{FILE_LABEL}.xlsx",
    index=False
)

time_by_category_month.to_excel(
    REPORTS_DIR / f"Monthly_Time_By_Category_{FILE_LABEL}.xlsx",
    index=False
)

time_by_tag_month.to_excel(
    REPORTS_DIR / f"Monthly_Time_By_Tag_{FILE_LABEL}.xlsx",
    index=False
)

# =====================================================
# CHARTS
# =====================================================

print("Creating charts...")

chart_daily = create_daily_trend_chart(
    daily_trend, "daily_trend.png"
)

chart_top_categories = create_bar_chart(
    top_categories_month,
    "Category", "Tickets",
    "Top Categories – Last 30 Days",
    "top_categories_month.png",
    top_n=10, horizontal=True
)

chart_top_tags = create_bar_chart(
    top_tags_month,
    "Tag", "Tickets",
    "Top Tags – Last 30 Days",
    "top_tags_month.png",
    top_n=10, horizontal=True
)

chart_open_categories = create_bar_chart(
    open_by_category,
    "Category", "Open Tickets",
    "Current Open Demand by Category",
    "open_categories_current.png",
    top_n=10, horizontal=True
)

chart_time_by_category = create_bar_chart(
    time_by_category_month,
    "Category", "TotalHours",
    "Time Logged by Category – Last 30 Days (Hours)",
    "time_by_category_month.png",
    top_n=10, horizontal=True
)

chart_time_by_tag = create_bar_chart(
    time_by_tag_month,
    "Tag", "TotalHours",
    "Time Logged by Tag – Last 30 Days (Hours)",
    "time_by_tag_month.png",
    top_n=10, horizontal=True
)

# =====================================================
# BUILD PRESENTATION
# =====================================================

print("Building PowerPoint...")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

# --- Slide 1: Title ---

slide = prs.slides.add_slide(blank_layout)

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(13.333), Inches(7.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = NAVY
shape.line.fill.background()

add_text(
    slide, REPORT_TITLE,
    0.8, 2.25, 11.8, 0.7,
    size=34, bold=True, color=WHITE
)
add_text(
    slide, REPORT_SUBTITLE,
    0.83, 3.05, 11.8, 0.4,
    size=18, color=WHITE
)
add_text(
    slide, "Eisenhower Center IT",
    0.83, 5.95, 11.8, 0.3,
    size=14, color=WHITE
)

# --- Slide 2: Executive Summary ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Executive Summary — Last 30 Days")
add_footer(slide)

created_count = len(created_this_month)
closed_count = len(closed_this_month)
open_count = len(open_now)
net_change = created_count - closed_count

if not top_categories_month.empty:
    top_cat_name = top_categories_month.iloc[0]["Category"]
    top_cat_count = top_categories_month.iloc[0]["Tickets"]
    top_cat_str = f"{top_cat_name} ({top_cat_count})"
else:
    top_cat_name = "N/A"
    top_cat_str = "N/A"

if not top_tags_month.empty:
    top_tag_name = top_tags_month.iloc[0]["Tag"]
    top_tag_count = top_tags_month.iloc[0]["Tickets"]
    top_tag_str = f"{top_tag_name} ({top_tag_count})"
else:
    top_tag_str = "N/A"

add_kpi_card(
    slide, "Tickets Created", created_count,
    0.45, 1.0, 2.0, 1.15
)
add_kpi_card(
    slide, "Tickets Closed", closed_count,
    2.60, 1.0, 2.0, 1.15,
    fill_color=GREEN_LIGHT
)
add_kpi_card(
    slide, "Net Change", f"{net_change:+d}",
    4.75, 1.0, 2.0, 1.15,
    fill_color=YELLOW_LIGHT if net_change > 0 else GREEN_LIGHT
)
add_kpi_card(
    slide, "Open Backlog", open_count,
    6.90, 1.0, 2.0, 1.15,
    fill_color=ORANGE_LIGHT
)
add_kpi_card(
    slide, "Time Logged (hrs)", total_hours_month,
    9.05, 1.0, 2.0, 1.15,
    fill_color=LIGHT_BLUE
)
add_kpi_card(
    slide, "Top Category", top_cat_name,
    11.20, 1.0, 1.9, 1.15,
    fill_color=LIGHT_BLUE
)

add_text(
    slide, "This Month's Story",
    0.55, 2.55, 12.5, 0.4,
    size=18, bold=True, color=NAVY
)

story_lines = [
    f"- {created_count} tickets came in, {closed_count} closed (net {net_change:+d}).",
    f"- Top demand category: {top_cat_str}.",
    f"- Top Zammad tag: {top_tag_str}.",
    f"- IT logged {total_hours_month} hours this month across {tickets_with_time_month} tickets ({time_coverage_month}% coverage).",
    f"- {open_count} tickets remain open, {aged_backlog} of which are over 30 days old.",
]

y = 3.05
for line in story_lines:
    add_text(
        slide, line,
        0.75, y, 12.2, 0.45,
        size=14, color=DARK
    )
    y += 0.55

# --- Slide 3: Daily Trend ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Daily Ticket Activity — Last 30 Days")
add_footer(slide)

add_chart_image(
    slide, chart_daily,
    0.55, 1.0, 12.2, 5.65
)

# --- Slide 4: Top Categories ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Highest Demand Categories — Last 30 Days")
add_footer(slide)

add_chart_image(
    slide, chart_top_categories,
    0.55, 1.0, 7.6, 5.7
)

add_table(
    slide, top_categories_month,
    8.45, 1.05, 4.3, 5.4,
    font_size=10, max_rows=12
)

# --- Slide 5: Top Tags ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Highest Demand Tags — Last 30 Days")
add_footer(slide)

add_chart_image(
    slide, chart_top_tags,
    0.55, 1.0, 7.6, 5.7
)

add_table(
    slide, top_tags_month,
    8.45, 1.05, 4.3, 5.4,
    font_size=10, max_rows=12
)

# --- Slide 6: Current Open ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Current Open Demand by Category")
add_footer(slide)

add_chart_image(
    slide, chart_open_categories,
    0.55, 1.0, 7.6, 5.7
)

add_table(
    slide, open_by_category,
    8.45, 1.05, 4.3, 5.4,
    font_size=10, max_rows=12
)

# --- Slide 7: Time by Category ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Time Logged by Category — Last 30 Days")
add_footer(slide)

add_chart_image(
    slide, chart_time_by_category,
    0.55, 1.0, 7.6, 5.7
)

time_cat_cols = [
    c for c in [
        "Category", "Tickets", "TicketsWithTime", "TotalHours"
    ]
    if c in time_by_category_month.columns
]

add_table(
    slide, time_by_category_month[time_cat_cols],
    8.35, 1.05, 4.65, 5.4,
    font_size=9, max_rows=12
)

# --- Slide 8: Time by Tag ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Time Logged by Tag — Last 30 Days")
add_footer(slide)

add_chart_image(
    slide, chart_time_by_tag,
    0.55, 1.0, 7.6, 5.7
)

time_tag_cols = [
    c for c in [
        "Tag", "Tickets", "TotalHours"
    ]
    if c in time_by_tag_month.columns
]

add_table(
    slide, time_by_tag_month[time_tag_cols],
    8.35, 1.05, 4.65, 5.4,
    font_size=9, max_rows=12
)

# --- Slide 9: Aged Backlog ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Aged Backlog (30+ Days)")
add_footer(slide)

if not high_risk_month.empty:
    aged_by_cat = (
        high_risk_month
        .groupby("Category")
        .agg(
            Tickets=("id", "count"),
            TotalMinutes=("time_unit", "sum")
        )
        .reset_index()
    )
    aged_by_cat["TotalHours"] = (
        (aged_by_cat["TotalMinutes"] / 60).round(1)
    )
    aged_by_cat = aged_by_cat.rename(
        columns={"Tickets": "30+ Day Tickets"}
    )
    aged_by_cat = aged_by_cat.sort_values(
        "30+ Day Tickets", ascending=False
    )

    aged_cols = [
        c for c in [
            "Category", "30+ Day Tickets", "TotalHours"
        ]
        if c in aged_by_cat.columns
    ]

    add_table(
        slide, aged_by_cat[aged_cols],
        0.65, 1.05, 5.3, 5.6,
        font_size=11, max_rows=15
    )

    detail_cols = [
        c for c in [
            "number", "title", "Category", "Days Open"
        ]
        if c in high_risk_month.columns
    ]

    if detail_cols:
        add_table(
            slide, high_risk_month[detail_cols],
            6.2, 1.05, 6.8, 5.6,
            font_size=8, max_rows=15
        )

else:
    add_text(
        slide,
        "No open tickets are over 30 days old.",
        0.75, 3.0, 12.0, 0.5,
        size=18, color=DARK, align=PP_ALIGN.CENTER
    )

# --- Slide 10: Recommendations ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Focus & Recommendations")
add_footer(slide)

recommendations = []

if net_change > 0:
    recommendations.append(
        f"Backlog is growing — {net_change} more tickets came in than closed this month."
    )
else:
    recommendations.append(
        f"Backlog is shrinking — closed {abs(net_change)} more tickets than created."
    )

if not top_categories_month.empty:
    recommendations.append(
        f"{top_cat_name} is the highest demand area this month — review documentation, training, or system fixes."
    )

if not time_by_category_month.empty:
    top_time_cat = time_by_category_month.iloc[0]["Category"]
    top_time_hours = time_by_category_month.iloc[0]["TotalHours"]
    recommendations.append(
        f"{top_time_cat} consumed the most logged time this month ({top_time_hours} hours)."
    )

if not top_tags_month.empty:
    recommendations.append(
        f"'{top_tag_name}' is the most-used tag — worth confirming if this reflects a real trend."
    )

if aged_backlog > 0:
    recommendations.append(
        f"{aged_backlog} tickets are over 30 days old — schedule weekly review."
    )

if time_coverage_month < 50:
    recommendations.append(
        f"Time logging coverage is {time_coverage_month}% — increase time tracking to make hours-based metrics reliable."
    )

recommendations.append(
    "Continue refining Zammad tags to improve future report fidelity."
)

y = 1.1
for idx, rec in enumerate(recommendations[:8], start=1):
    add_text(
        slide, f"{idx}. {rec}",
        0.75, y, 11.9, 0.6,
        size=15, color=DARK
    )
    y += 0.7

# =====================================================
# SAVE
# =====================================================

prs.save(PPTX_FILE)

print()
print("=" * 60)
print("MONTHLY REPORT CREATED")
print("=" * 60)
print(PPTX_FILE)
print()
print("Done.")