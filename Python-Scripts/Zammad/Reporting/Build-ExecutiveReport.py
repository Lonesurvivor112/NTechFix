import pandas as pd
import matplotlib.pyplot as plt
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# =====================================================
# CONFIG
# =====================================================

OUTPUT_DIR = Path("output")

REPORTS_DIR = Path(
    r"C:\Users\nfalzetti\eisenhowercenter.com"
    r"\IT - Documents\Scripts\Nick Scripts"
    r"\Python\Zammad py\ZammadReports\Reports"
)
REPORTS_DIR.mkdir(parents=True, exist_ok=True)

CHART_DIR = REPORTS_DIR / "executive_charts"
CHART_DIR.mkdir(exist_ok=True)

REPORT_DATE = pd.Timestamp.utcnow().tz_localize(None)
FILE_LABEL = REPORT_DATE.strftime("%Y-%m-%d")

PPTX_FILE = (
    REPORTS_DIR /
    f"EC_IT_Executive_Report_{FILE_LABEL}.pptx"
)

KPI_FILE = OUTPUT_DIR / "KPI_Summary.xlsx"
MONTHLY_FILE = OUTPUT_DIR / "Monthly_Trend.xlsx"
MONTHLY_TIME_FILE = OUTPUT_DIR / "Monthly_TimeSpent.xlsx"
TOP_CATEGORIES_FILE = OUTPUT_DIR / "Top_Categories.xlsx"
OPEN_CATEGORIES_FILE = OUTPUT_DIR / "Open_Categories.xlsx"
BACKLOG_AGING_FILE = OUTPUT_DIR / "Backlog_Aging.xlsx"
HIGH_RISK_FILE = OUTPUT_DIR / "High_Risk_Backlog.xlsx"
HIGH_RISK_BY_CATEGORY_FILE = OUTPUT_DIR / "High_Risk_By_Category.xlsx"
REPEAT_FILE = OUTPUT_DIR / "Repeat_Problems.xlsx"
CATEGORY_SUMMARY_FILE = OUTPUT_DIR / "Category_Summary.xlsx"
OPEN_TIME_FILE = OUTPUT_DIR / "Open_TimeSpent.xlsx"

REPORT_TITLE = "IT Service Desk Executive Report"
REPORT_SUBTITLE = "All-Time Overview"

# =====================================================
# STYLE
# =====================================================

NAVY = RGBColor(31, 78, 121)
DARK = RGBColor(38, 38, 38)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(221, 235, 247)
LIGHT_GRAY = RGBColor(242, 242, 242)
GREEN_LIGHT = RGBColor(226, 239, 218)
YELLOW_LIGHT = RGBColor(255, 242, 204)
ORANGE_LIGHT = RGBColor(252, 228, 214)
RED_LIGHT = RGBColor(255, 199, 206)
GRAY_LINE = RGBColor(210, 210, 210)
GRAY_TEXT = RGBColor(120, 120, 120)

FONT = "Aptos"

# =====================================================
# HELPERS
# =====================================================

def read_excel_safe(path):
    if not path.exists():
        print(f"Missing file: {path}")
        return pd.DataFrame()

    return pd.read_excel(path)


def get_kpi_dict(kpi_df):
    if kpi_df.empty:
        return {}

    return dict(
        zip(
            kpi_df["Metric"],
            kpi_df["Value"]
        )
    )


def add_header(slide, title):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.62)
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    textbox = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.1),
        Inches(12.6), Inches(0.4)
    )

    frame = textbox.text_frame
    frame.clear()

    p = frame.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_footer(slide):
    textbox = slide.shapes.add_textbox(
        Inches(0.35), Inches(7.16),
        Inches(12.6), Inches(0.22)
    )

    frame = textbox.text_frame
    p = frame.paragraphs[0]
    p.text = (
        f"Eisenhower Center IT | "
        f"Executive Report | "
        f"{REPORT_DATE.strftime('%b %d, %Y')}"
    )
    p.font.name = FONT
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.RIGHT


def add_text(
    slide, text, x, y, w, h,
    size=16, bold=False, color=DARK, align=None
):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    p = frame.paragraphs[0]
    p.text = str(text)
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color

    if align is not None:
        p.alignment = align

    return box


def add_kpi_card(
    slide, label, value, x, y, w, h,
    fill_color=LIGHT_BLUE
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = GRAY_LINE

    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)

    p1 = frame.paragraphs[0]
    p1.text = str(value)
    p1.font.name = FONT
    p1.font.size = Pt(26)
    p1.font.bold = True
    p1.font.color.rgb = NAVY
    p1.alignment = PP_ALIGN.CENTER

    p2 = frame.add_paragraph()
    p2.text = str(label)
    p2.font.name = FONT
    p2.font.size = Pt(10)
    p2.font.color.rgb = DARK
    p2.alignment = PP_ALIGN.CENTER


def add_table(
    slide, dataframe, x, y, w, h,
    font_size=9, max_rows=10
):
    if dataframe.empty:
        add_text(
            slide, "No data available.",
            x, y, w, h, size=12
        )
        return

    df = dataframe.head(max_rows).copy()

    rows = len(df) + 1
    cols = len(df.columns)

    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(x), Inches(y),
        Inches(w), Inches(h)
    )

    table = table_shape.table

    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = FONT
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = WHITE

    for row_idx, (_, row) in enumerate(
        df.iterrows(), start=1
    ):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)

            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = DARK


def create_bar_chart(
    df, label_col, value_col, title, filename,
    top_n=10, horizontal=True
):
    chart_path = CHART_DIR / filename

    if df.empty:
        return None

    if label_col not in df.columns or value_col not in df.columns:
        return None

    data = df.head(top_n).copy()

    plt.figure(figsize=(9, 4.8))

    if horizontal:
        data = data.sort_values(
            value_col, ascending=True
        )
        plt.barh(
            data[label_col].astype(str),
            data[value_col]
        )
        plt.xlabel(value_col)
    else:
        plt.bar(
            data[label_col].astype(str),
            data[value_col]
        )
        plt.ylabel(value_col)
        plt.xticks(rotation=35, ha="right")

    plt.title(title)
    plt.tight_layout()
    plt.savefig(chart_path, dpi=160)
    plt.close()

    return chart_path


def create_monthly_trend_chart(df, filename):
    chart_path = CHART_DIR / filename

    if df.empty:
        return None

    required = {"Month", "Created", "Closed"}

    if not required.issubset(df.columns):
        return None

    data = df.tail(18).copy()

    plt.figure(figsize=(10, 4.8))

    plt.plot(
        data["Month"], data["Created"],
        marker="o", label="Created"
    )
    plt.plot(
        data["Month"], data["Closed"],
        marker="o", label="Closed"
    )

    if "Net" in data.columns:
        plt.bar(
            data["Month"], data["Net"],
            alpha=0.25, label="Net"
        )

    plt.title("Ticket Created vs Closed Trend")
    plt.xlabel("Month")
    plt.ylabel("Tickets")
    plt.xticks(rotation=45, ha="right")
    plt.legend()
    plt.tight_layout()
    plt.savefig(chart_path, dpi=160)
    plt.close()

    return chart_path


def create_monthly_time_chart(df, filename):
    chart_path = CHART_DIR / filename

    if df.empty:
        return None

    if "Month" not in df.columns or "TotalHours" not in df.columns:
        return None

    data = df.tail(18).copy()

    plt.figure(figsize=(10, 4.8))

    plt.bar(
        data["Month"],
        data["TotalHours"],
        color="#2E75B6"
    )

    plt.title("Monthly Time Logged (Hours)")
    plt.xlabel("Month")
    plt.ylabel("Hours Logged")
    plt.xticks(rotation=45, ha="right")
    plt.tight_layout()
    plt.savefig(chart_path, dpi=160)
    plt.close()

    return chart_path


def add_chart_image(slide, chart_path, x, y, w, h):
    if chart_path and Path(chart_path).exists():
        slide.shapes.add_picture(
            str(chart_path),
            Inches(x), Inches(y),
            Inches(w), Inches(h)
        )
    else:
        add_text(
            slide,
            "Chart could not be generated.",
            x, y, w, h, size=12
        )


def metric_value(kpis, name, default="N/A"):
    return kpis.get(name, default)


def build_insights(kpis, open_categories, category_summary, monthly):
    insights = []

    open_backlog = kpis.get("Open Backlog", 0)
    aged_backlog = kpis.get("30+ Day Backlog", 0)
    aged_pct = kpis.get("30+ Day Backlog %", None)

    try:
        open_backlog_int = int(open_backlog)
        aged_backlog_int = int(aged_backlog)
    except:
        open_backlog_int = 0
        aged_backlog_int = 0

    if aged_pct is not None:
        insights.append(
            f"{aged_pct}% of the open backlog is older than 30 days "
            f"({aged_backlog_int} of {open_backlog_int} open tickets)."
        )

    if not open_categories.empty:
        row = open_categories.iloc[0]
        insights.append(
            f"The largest current open category is {row['Category']} "
            f"with {row['Open Tickets']} open tickets."
        )

    if not category_summary.empty:
        row = category_summary.iloc[0]
        hours = row.get("TotalHours", "N/A")
        insights.append(
            f"The most time-consuming category is {row['Category']} "
            f"with {hours} hours logged."
        )

    top_time_cat = kpis.get("Top Category by Time")
    if top_time_cat:
        insights.append(
            f"Top category by time: {top_time_cat}."
        )

    top_time_tag = kpis.get("Top Tag by Time")
    if top_time_tag:
        insights.append(
            f"Top tag by time: {top_time_tag}."
        )

    if not monthly.empty and "Net" in monthly.columns:
        recent = monthly.iloc[-1]
        insights.append(
            f"Latest month net: {recent['Net']} tickets "
            f"({recent['Created']} created, {recent['Closed']} closed)."
        )

    time_coverage = kpis.get("Time Logging Coverage %")
    if time_coverage is not None:
        insights.append(
            f"Time logging coverage: {time_coverage}% "
            f"of tickets have logged minutes."
        )

    if len(insights) == 0:
        insights.append(
            "No automated insights could be generated."
        )

    return insights


# =====================================================
# LOAD DATA
# =====================================================

print("Loading data...")

kpi_df = read_excel_safe(KPI_FILE)
monthly = read_excel_safe(MONTHLY_FILE)
monthly_time = read_excel_safe(MONTHLY_TIME_FILE)
top_categories = read_excel_safe(TOP_CATEGORIES_FILE)
open_categories = read_excel_safe(OPEN_CATEGORIES_FILE)
backlog_aging = read_excel_safe(BACKLOG_AGING_FILE)
high_risk = read_excel_safe(HIGH_RISK_FILE)
high_risk_by_category = read_excel_safe(HIGH_RISK_BY_CATEGORY_FILE)
repeat_problems = read_excel_safe(REPEAT_FILE)
category_summary = read_excel_safe(CATEGORY_SUMMARY_FILE)
open_time = read_excel_safe(OPEN_TIME_FILE)

kpis = get_kpi_dict(kpi_df)

# =====================================================
# CREATE CHARTS
# =====================================================

print("Creating charts...")

monthly_chart = create_monthly_trend_chart(
    monthly, "monthly_trend.png"
)

monthly_time_chart = create_monthly_time_chart(
    monthly_time, "monthly_time.png"
)

open_categories_chart = create_bar_chart(
    open_categories,
    "Category", "Open Tickets",
    "Open Tickets by Category",
    "open_categories.png",
    top_n=10, horizontal=True
)

backlog_chart = create_bar_chart(
    backlog_aging,
    "Bucket", "Open Tickets",
    "Backlog Aging",
    "backlog_aging.png",
    top_n=10, horizontal=False
)

top_categories_chart = create_bar_chart(
    top_categories,
    "Category", "Tickets",
    "Historical Tickets by Category",
    "top_categories.png",
    top_n=10, horizontal=True
)

time_by_category_chart = create_bar_chart(
    category_summary,
    "Category", "TotalHours",
    "Time Logged by Category (Hours)",
    "time_by_category.png",
    top_n=10, horizontal=True
)

open_time_chart = create_bar_chart(
    open_time,
    "Category", "TotalHours",
    "Time Logged on Open Tickets (Hours)",
    "open_time.png",
    top_n=10, horizontal=True
)

# =====================================================
# BUILD PRESENTATION
# =====================================================

print("Building PowerPoint...")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

# --- Slide 1: Title ---

slide = prs.slides.add_slide(blank_layout)

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(13.333), Inches(7.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = NAVY
shape.line.fill.background()

add_text(
    slide, REPORT_TITLE,
    0.8, 2.25, 11.8, 0.7,
    size=34, bold=True, color=WHITE
)
add_text(
    slide, REPORT_SUBTITLE,
    0.83, 3.05, 11.8, 0.4,
    size=18, color=WHITE
)
add_text(
    slide, f"Generated {REPORT_DATE.strftime('%B %d, %Y')}",
    0.83, 3.65, 11.8, 0.4,
    size=14, color=WHITE
)
add_text(
    slide, "Eisenhower Center IT",
    0.83, 5.95, 11.8, 0.3,
    size=14, color=WHITE
)

# --- Slide 2: Executive Summary ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Executive Summary")
add_footer(slide)

add_kpi_card(
    slide, "Tickets Created",
    metric_value(kpis, "Tickets Created"),
    0.45, 1.0, 1.9, 1.1
)
add_kpi_card(
    slide, "Tickets Closed",
    metric_value(kpis, "Tickets Closed"),
    2.55, 1.0, 1.9, 1.1,
    fill_color=GREEN_LIGHT
)
add_kpi_card(
    slide, "Open Backlog",
    metric_value(kpis, "Open Backlog"),
    4.65, 1.0, 1.9, 1.1,
    fill_color=YELLOW_LIGHT
)
add_kpi_card(
    slide, "30+ Day Backlog",
    metric_value(kpis, "30+ Day Backlog"),
    6.75, 1.0, 1.9, 1.1,
    fill_color=ORANGE_LIGHT
)
add_kpi_card(
    slide, "Time Logged (hrs)",
    metric_value(kpis, "Total Time Logged (hrs)"),
    8.85, 1.0, 1.9, 1.1,
    fill_color=LIGHT_BLUE
)
add_kpi_card(
    slide, "Time Coverage %",
    metric_value(kpis, "Time Logging Coverage %"),
    10.95, 1.0, 1.9, 1.1,
    fill_color=RED_LIGHT
)

add_text(
    slide, "Automated Insights",
    0.55, 2.55, 5.8, 0.35,
    size=18, bold=True, color=NAVY
)

insights = build_insights(
    kpis, open_categories, category_summary, monthly
)

y = 3.0
for insight in insights[:6]:
    add_text(
        slide, f"- {insight}",
        0.75, y, 5.85, 0.55,
        size=12, color=DARK
    )
    y += 0.55

add_text(
    slide, "Management Focus",
    7.0, 2.55, 5.8, 0.35,
    size=18, bold=True, color=NAVY
)

focus_items = [
    "Reduce aged ticket backlog.",
    "Focus on highest-time categories for process improvement.",
    "Continue tagging discipline in Zammad.",
    "Log time on more tickets to improve coverage.",
    "Use recurring issues to guide documentation and training.",
    "Track created vs closed monthly."
]

y = 3.0
for item in focus_items:
    add_text(
        slide, f"- {item}",
        7.2, y, 5.8, 0.5,
        size=12, color=DARK
    )
    y += 0.55

# --- Slide 3: Monthly Trend ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Monthly Ticket Trend")
add_footer(slide)

add_chart_image(
    slide, monthly_chart,
    0.55, 1.0, 12.2, 5.65
)

# --- Slide 4: Monthly Time Logged ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Monthly Time Logged")
add_footer(slide)

add_chart_image(
    slide, monthly_time_chart,
    0.55, 1.0, 12.2, 5.65
)

# --- Slide 5: Backlog Aging ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Backlog Aging")
add_footer(slide)

add_chart_image(
    slide, backlog_chart,
    0.55, 1.0, 7.2, 5.6
)

add_text(
    slide, "Why This Matters",
    8.1, 1.05, 4.5, 0.35,
    size=18, bold=True, color=NAVY
)

aged_pct = metric_value(kpis, "30+ Day Backlog %")
aged_count = metric_value(kpis, "30+ Day Backlog")
open_count = metric_value(kpis, "Open Backlog")

backlog_text = (
    f"The current backlog contains {open_count} open tickets. "
    f"{aged_count} are over 30 days old, representing {aged_pct}% "
    f"of the open backlog."
)

add_text(
    slide, backlog_text,
    8.1, 1.55, 4.55, 1.4,
    size=14, color=DARK
)

add_text(
    slide,
    "Recommended focus: review high-risk tickets weekly and close or reclassify stale items.",
    8.1, 3.25, 4.55, 1.4,
    size=14, color=DARK
)

# --- Slide 6: Open Categories ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Current Open Demand by Category")
add_footer(slide)

add_chart_image(
    slide, open_categories_chart,
    0.55, 1.0, 7.6, 5.7
)

add_table(
    slide, open_categories,
    8.45, 1.05, 4.3, 5.4,
    font_size=9, max_rows=12
)

# --- Slide 7: Historical Categories ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Historical Ticket Demand by Category")
add_footer(slide)

add_chart_image(
    slide, top_categories_chart,
    0.55, 1.0, 7.6, 5.7
)

add_table(
    slide, top_categories,
    8.45, 1.05, 4.3, 5.4,
    font_size=9, max_rows=12
)

# --- Slide 8: Time by Category ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Time Logged by Category")
add_footer(slide)

add_chart_image(
    slide, time_by_category_chart,
    0.55, 1.0, 7.6, 5.7
)

time_cols = [
    c for c in [
        "Category", "Tickets",
        "TicketsWithTime", "TotalHours",
        "AvgMinutesPerTicket", "TimeCoverage %"
    ]
    if c in category_summary.columns
]

if time_cols:
    add_table(
        slide, category_summary[time_cols],
        8.35, 1.05, 4.65, 5.4,
        font_size=7, max_rows=12
    )

# --- Slide 9: Open Time by Category ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Time Logged on Currently Open Tickets")
add_footer(slide)

add_chart_image(
    slide, open_time_chart,
    0.55, 1.0, 7.6, 5.7
)

open_time_cols = [
    c for c in [
        "Category", "OpenTickets",
        "TicketsWithTime", "TotalHours",
        "AvgMinutesPerTicket"
    ]
    if c in open_time.columns
]

if open_time_cols:
    add_table(
        slide, open_time[open_time_cols],
        8.35, 1.05, 4.65, 5.4,
        font_size=8, max_rows=12
    )

# --- Slide 10: High Risk Backlog ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "High-Risk Backlog")
add_footer(slide)

add_text(
    slide, "Open tickets over 30 days old by category.",
    0.55, 0.85, 12.0, 0.35,
    size=13, color=DARK
)

hr_cat_cols = [
    c for c in [
        "Category", "30+ Day Tickets",
        "TotalMinutes", "TotalHours"
    ]
    if c in high_risk_by_category.columns
]

add_table(
    slide, high_risk_by_category[hr_cat_cols],
    0.65, 1.35, 5.3, 5.4,
    font_size=10, max_rows=12
)

high_risk_display_cols = [
    c for c in [
        "number", "title", "Category", "Days Open"
    ]
    if c in high_risk.columns
]

if high_risk_display_cols:
    add_table(
        slide, high_risk[high_risk_display_cols],
        6.2, 1.35, 6.75, 5.4,
        font_size=7, max_rows=12
    )

# --- Slide 11: Recurring Issues ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Recurring Issues")
add_footer(slide)

add_text(
    slide,
    "Normalized repeat ticket titles help identify documentation, training, or system improvement opportunities.",
    0.55, 0.9, 12.2, 0.45,
    size=13, color=DARK
)

add_table(
    slide, repeat_problems,
    0.75, 1.45, 11.85, 5.4,
    font_size=10, max_rows=15
)

# --- Slide 12: Recommendations ---

slide = prs.slides.add_slide(blank_layout)
add_header(slide, "Operational Recommendations")
add_footer(slide)

recommendations = []

aged_pct_metric = metric_value(kpis, "30+ Day Backlog %", 0)

try:
    aged_pct_number = float(aged_pct_metric)

    if aged_pct_number >= 50:
        recommendations.append(
            "Create a weekly aged-ticket review focused on tickets older than 30 days."
        )
except:
    pass

if not open_categories.empty:
    first_category = str(open_categories.iloc[0]["Category"])
    recommendations.append(
        f"Review the {first_category} category — it has the highest current open ticket volume."
    )

if not category_summary.empty:
    top_time_cat = str(category_summary.iloc[0]["Category"])
    recommendations.append(
        f"Investigate {top_time_cat} process improvements — it consumes the most logged IT time."
    )

if not top_categories.empty:
    if "Other" in top_categories["Category"].astype(str).values:
        recommendations.append(
            "Continue refining Zammad tags to reduce tickets categorized as 'Other'."
        )

time_cov = metric_value(kpis, "Time Logging Coverage %", 0)

try:
    if float(time_cov) < 50:
        recommendations.append(
            f"Time logging coverage is {time_cov}% — improve time tracking discipline to make time metrics more reliable."
        )
except:
    pass

recommendations.extend([
    "Use recurring issues as input for quick staff guides and SOPs.",
    "Track monthly Net tickets to see whether backlog is growing or shrinking."
])

y = 1.1
for idx, rec in enumerate(recommendations[:8], start=1):
    add_text(
        slide, f"{idx}. {rec}",
        0.75, y, 11.9, 0.6,
        size=15, color=DARK
    )
    y += 0.65

# =====================================================
# SAVE
# =====================================================

prs.save(PPTX_FILE)

print()
print("=" * 60)
print("EXECUTIVE REPORT CREATED")
print("=" * 60)
print(PPTX_FILE)
print()
print("Done.")