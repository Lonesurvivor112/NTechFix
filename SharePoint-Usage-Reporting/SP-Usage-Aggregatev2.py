r"""
Aggregate SharePoint Site Analytics exports into an executive Excel dashboard
and (optionally) a management-ready PowerPoint deck.

Usage:
    python SP-Usage-Aggregatev3.py <input_folder> [output_xlsx] [--pptx output_pptx]

Example:
    python SP-Usage-Aggregatev3.py "C:\Reports\downloads" "C:\Reports\dashboard.xlsx" --pptx "C:\Reports\briefing.pptx"

Focus: 30-Day Visits and 30-Day Unique Viewers, with granular reporting of
lower-traffic sites and a callout for sites with no 30-day activity.
"""
import sys
import argparse
import re
from pathlib import Path
from datetime import datetime
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.chart.label import DataLabelList
from openpyxl.utils import get_column_letter
from openpyxl.formatting.rule import ColorScaleRule


def force_axes_visible(chart, x_format='General', y_format='#,##0'):
    """Force tick labels on both axes to display, with readable formatting."""
    chart.x_axis.delete = False
    chart.x_axis.tickLblPos = 'low'
    chart.x_axis.majorTickMark = 'out'
    chart.x_axis.minorTickMark = 'none'
    chart.x_axis.number_format = x_format
    chart.y_axis.delete = False
    chart.y_axis.tickLblPos = 'nextTo'
    chart.y_axis.majorTickMark = 'out'
    chart.y_axis.minorTickMark = 'none'
    chart.y_axis.number_format = y_format
    chart.y_axis.crosses = 'autoZero'
    from openpyxl.chart.axis import ChartLines
    chart.y_axis.majorGridlines = ChartLines()


# ---------- Parsing ----------

def parse_overall_traffic(df, site_name):
    """Pull aggregate KPIs and the daily time series from 'Overall traffic'."""
    out = {
        'site': site_name,
        'viewers_7d': None, 'visits_7d': None,
        'viewers_30d': None, 'visits_30d': None,
        'viewers_90d': None, 'visits_90d': None,
        'viewers_all': None, 'visits_all': None,
        'daily': pd.DataFrame(columns=['date', 'unique_viewers', 'site_visits']),
    }

    arr = df.values
    agg_header = None
    daily_header = None
    for i, row in enumerate(arr):
        cells = [str(c) if pd.notna(c) else '' for c in row]
        if cells[0] == 'Duration' and 'Unique viewers' in cells[1:]:
            agg_header = i
        if cells[0] == 'Date' and 'Unique viewers' in cells[1:]:
            daily_header = i
            break

    def to_int(v):
        try:
            if v is None or pd.isna(v):
                return None
            if isinstance(v, str) and not v.replace('.', '', 1).isdigit():
                return None
            return int(float(v))
        except Exception:
            return None

    if agg_header is not None:
        for r in range(agg_header + 1, agg_header + 6):
            if r >= len(arr):
                break
            label = str(arr[r, 0]) if pd.notna(arr[r, 0]) else ''
            v = to_int(arr[r, 1])
            s = to_int(arr[r, 2])
            if 'Last 7' in label:
                out['viewers_7d'], out['visits_7d'] = v, s
            elif 'Last 30' in label:
                out['viewers_30d'], out['visits_30d'] = v, s
            elif 'Last 90' in label:
                out['viewers_90d'], out['visits_90d'] = v, s
            elif 'All time' in label:
                out['viewers_all'], out['visits_all'] = v, s

    if daily_header is not None:
        rows = []
        for r in range(daily_header + 1, len(arr)):
            d = arr[r, 0]
            if pd.isna(d):
                continue
            try:
                date = pd.to_datetime(d).date()
            except Exception:
                continue
            rows.append((date, to_int(arr[r, 1]) or 0, to_int(arr[r, 2]) or 0))
        if rows:
            out['daily'] = pd.DataFrame(rows, columns=['date', 'unique_viewers', 'site_visits'])

    return out


def parse_popular_content(df, site_name):
    """Top content items from the 'Popular content' sheet."""
    arr = df.values
    header_row = None
    for i, row in enumerate(arr):
        cells = [str(c) if pd.notna(c) else '' for c in row]
        if cells[0] == 'Content' and any('viewers' in c.lower() for c in cells):
            header_row = i
            break
    if header_row is None:
        return pd.DataFrame()

    rows = []
    for r in range(header_row + 1, len(arr)):
        name = arr[r, 0]
        if pd.isna(name) or 'Click here' in str(name):
            continue
        rows.append({
            'site': site_name,
            'content': str(name),
            'type': str(arr[r, 1]) if pd.notna(arr[r, 1]) else '',
            'viewers_7d': arr[r, 2] if pd.notna(arr[r, 2]) else 0,
            'visits_7d': arr[r, 3] if pd.notna(arr[r, 3]) else 0,
        })
    return pd.DataFrame(rows)


def parse_device_usage(df, site_name):
    """Sum device-type visits across the 90-day window."""
    arr = df.values
    header_row = None
    for i, row in enumerate(arr):
        cells = [str(c) if pd.notna(c) else '' for c in row]
        if cells[0] == 'Date' and any('Desktop' in c for c in cells):
            header_row = i
            break
    if header_row is None:
        return {'desktop': 0, 'mobile_app': 0, 'mobile_web': 0, 'tablet': 0, 'other': 0}

    totals = {'desktop': 0, 'mobile_app': 0, 'mobile_web': 0, 'tablet': 0, 'other': 0}
    for r in range(header_row + 1, len(arr)):
        if pd.isna(arr[r, 0]):
            continue
        try:
            totals['desktop'] += int(arr[r, 1] or 0)
            totals['mobile_app'] += int(arr[r, 2] or 0)
            totals['mobile_web'] += int(arr[r, 3] or 0)
            totals['tablet'] += int(arr[r, 4] or 0)
            totals['other'] += int(arr[r, 5] or 0)
        except (ValueError, TypeError):
            continue
    return totals


def site_name_from_workbook(sheets):
    """Extract the site name from cell A1 of any sheet."""
    for df in sheets.values():
        if df.shape[0] > 0 and df.shape[1] > 0:
            first = str(df.columns[0])
            m = re.match(r'Site name:\s*(.+)', first)
            if m:
                return m.group(1).strip()
    return None


def load_site(path):
    """Parse one xlsx export. Returns dict with all extracted data, or None on failure."""
    try:
        sheets = pd.read_excel(path, sheet_name=None)
    except Exception as e:
        print(f"  ! could not read {path.name}: {e}")
        return None

    site = site_name_from_workbook(sheets) or path.stem

    overall = parse_overall_traffic(sheets.get('Overall traffic', pd.DataFrame()), site)
    popular = parse_popular_content(sheets.get('Popular content', pd.DataFrame()), site)
    devices = parse_device_usage(sheets.get('Usage by device', pd.DataFrame()), site)

    return {
        'site': site,
        'file': path.name,
        'overall': overall,
        'popular': popular,
        'devices': devices,
    }


# ---------- Output ----------

HEADER_FILL = PatternFill('solid', start_color='1F4E78')
HEADER_FONT = Font(bold=True, color='FFFFFF', name='Calibri', size=11)
TITLE_FONT = Font(bold=True, size=16, color='1F4E78', name='Calibri')
SUBTITLE_FONT = Font(italic=True, size=10, color='595959', name='Calibri')
KPI_LABEL_FONT = Font(bold=True, size=10, color='595959', name='Calibri')
KPI_VALUE_FONT = Font(bold=True, size=20, color='1F4E78', name='Calibri')
THIN_BORDER = Border(
    left=Side(style='thin', color='D0D0D0'),
    right=Side(style='thin', color='D0D0D0'),
    top=Side(style='thin', color='D0D0D0'),
    bottom=Side(style='thin', color='D0D0D0'),
)


def style_header_row(ws, row, ncols):
    for c in range(1, ncols + 1):
        cell = ws.cell(row=row, column=c)
        cell.fill = HEADER_FILL
        cell.font = HEADER_FONT
        cell.alignment = Alignment(horizontal='center', vertical='center')
        cell.border = THIN_BORDER


def autosize_columns(ws, widths):
    for col, w in widths.items():
        ws.column_dimensions[col].width = w


def write_summary_sheet(wb, sites_data):
    """Executive Summary: KPIs + ranked site table + top-10 & lower-traffic charts + inactive callout."""
    ws = wb.create_sheet('Executive Summary', 0)
    ws.sheet_view.showGridLines = False

    ws['A1'] = 'SharePoint Site Usage — Executive Summary'
    ws['A1'].font = TITLE_FONT
    ws.merge_cells('A1:L1')
    ws['A2'] = (f'Generated {datetime.now().strftime("%B %d, %Y")}'
                f' • {len(sites_data)} sites • Focus: last 30 days')
    ws['A2'].font = SUBTITLE_FONT
    ws.merge_cells('A2:L2')

    total_viewers_30d = sum((s['overall']['viewers_30d'] or 0) for s in sites_data)
    total_visits_30d = sum((s['overall']['visits_30d'] or 0) for s in sites_data)
    total_viewers_7d = sum((s['overall']['viewers_7d'] or 0) for s in sites_data)
    total_visits_7d = sum((s['overall']['visits_7d'] or 0) for s in sites_data)
    active_sites = sum(1 for s in sites_data if (s['overall']['visits_30d'] or 0) > 0)
    inactive_sites = len(sites_data) - active_sites

    HIGHLIGHT_FILL = PatternFill('solid', start_color='DDEBF7')
    SECONDARY_FILL = PatternFill('solid', start_color='F2F2F2')
    HIGHLIGHT_VALUE_FONT = Font(bold=True, size=24, color='1F4E78', name='Calibri')

    kpi_cards = [
        ('A', 'B', '30-Day Visits',         total_visits_30d,   HIGHLIGHT_FILL, HIGHLIGHT_VALUE_FONT, 'Total site visits • last 30d'),
        ('C', 'D', '30-Day Unique Viewers', total_viewers_30d,  HIGHLIGHT_FILL, HIGHLIGHT_VALUE_FONT, 'Distinct users • last 30d'),
        ('E', 'F', 'Active Sites (30d)',    f'{active_sites} / {len(sites_data)}', SECONDARY_FILL, KPI_VALUE_FONT, 'Sites with 30d activity'),
        ('G', 'H', 'Inactive Sites (30d)',  inactive_sites,     SECONDARY_FILL, KPI_VALUE_FONT, 'Zero visits in last 30d'),
        ('I', 'J', '7-Day Visits',          total_visits_7d,    SECONDARY_FILL, KPI_VALUE_FONT, 'Trailing week total'),
        ('K', 'L', '7-Day Unique Viewers',  total_viewers_7d,   SECONDARY_FILL, KPI_VALUE_FONT, 'Trailing week users'),
    ]

    kpi_top = 4
    for start, end, label, value, fill, vfont, note in kpi_cards:
        ws.merge_cells(f'{start}{kpi_top}:{end}{kpi_top}')
        c = ws[f'{start}{kpi_top}']
        c.value = label.upper()
        c.font = Font(bold=True, size=9, color='595959', name='Calibri')
        c.alignment = Alignment(horizontal='center', vertical='center')
        c.fill = fill

        ws.merge_cells(f'{start}{kpi_top+1}:{end}{kpi_top+1}')
        v = ws[f'{start}{kpi_top+1}']
        v.value = value
        v.font = vfont
        v.alignment = Alignment(horizontal='center', vertical='center')
        v.fill = fill
        if isinstance(value, (int, float)):
            v.number_format = '#,##0'

        ws.merge_cells(f'{start}{kpi_top+2}:{end}{kpi_top+2}')
        n = ws[f'{start}{kpi_top+2}']
        n.value = note
        n.font = Font(italic=True, size=8, color='7F7F7F', name='Calibri')
        n.alignment = Alignment(horizontal='center', vertical='center')
        n.fill = fill

    ws.row_dimensions[kpi_top].height = 18
    ws.row_dimensions[kpi_top+1].height = 34
    ws.row_dimensions[kpi_top+2].height = 16

    # Ranked table (30d Visits + 30d Viewers front and center)
    table_start = 9
    headers = ['Site', 'Visits 30d', 'Viewers 30d', 'Visits 7d', 'Viewers 7d',
               'Visits 90d', 'Visits All-time', 'Status']
    for i, h in enumerate(headers, 1):
        ws.cell(row=table_start, column=i, value=h)
    style_header_row(ws, table_start, len(headers))

    rows = []
    for s in sites_data:
        o = s['overall']
        rows.append([
            s['site'],
            o['visits_30d'] or 0,
            o['viewers_30d'] or 0,
            o['visits_7d'] or 0,
            o['viewers_7d'] or 0,
            o['visits_90d'] or 0,
            o['visits_all'] or 0,
        ])
    rows.sort(key=lambda r: r[1], reverse=True)

    for i, r in enumerate(rows):
        excel_row = table_start + 1 + i
        for j, val in enumerate(r):
            cell = ws.cell(row=excel_row, column=j + 1, value=val)
            if j > 0:
                cell.number_format = '#,##0'
                cell.alignment = Alignment(horizontal='right')
            cell.border = THIN_BORDER
        status_col = len(headers)
        st = ws.cell(row=excel_row, column=status_col,
                     value=f'=IF(B{excel_row}=0,"Inactive",IF(B{excel_row}<10,"Low",IF(B{excel_row}<100,"Moderate","Active")))')
        st.alignment = Alignment(horizontal='center')
        st.border = THIN_BORDER
        if i % 2 == 1:
            for j in range(len(headers)):
                ws.cell(row=excel_row, column=j + 1).fill = PatternFill('solid', start_color='F7F9FC')

    last_row = table_start + len(rows)

    total_row = last_row + 1
    tot_label = ws.cell(row=total_row, column=1, value='TOTAL')
    tot_label.font = Font(bold=True, color='FFFFFF')
    tot_label.fill = PatternFill('solid', start_color='1F4E78')
    tot_label.alignment = Alignment(horizontal='right')
    for col_idx in range(2, 8):
        col_letter = get_column_letter(col_idx)
        cell = ws.cell(row=total_row, column=col_idx,
                       value=f'=SUM({col_letter}{table_start + 1}:{col_letter}{last_row})')
        cell.font = Font(bold=True, color='FFFFFF')
        cell.number_format = '#,##0'
        cell.fill = PatternFill('solid', start_color='1F4E78')
        cell.alignment = Alignment(horizontal='right')

    if len(rows) > 0:
        for col_letter in ('B', 'C'):
            ws.conditional_formatting.add(
                f'{col_letter}{table_start + 1}:{col_letter}{last_row}',
                ColorScaleRule(start_type='min', start_color='FFFFFF',
                               mid_type='percentile', mid_value=50, mid_color='9CC2E5',
                               end_type='max', end_color='1F4E78'),
            )

    autosize_columns(ws, {
        'A': 40, 'B': 13, 'C': 14, 'D': 12, 'E': 12,
        'F': 13, 'G': 17, 'H': 12,
        'I': 13, 'J': 14, 'K': 13, 'L': 14,
    })
    ws.freeze_panes = 'B10'

    # Top-10 chart
    chart_col = 14
    chart_label_col = chart_col
    chart_visits_col = chart_col + 1
    chart_viewers_col = chart_col + 2

    ws.cell(row=table_start, column=chart_label_col, value='Site (Top 10)')
    ws.cell(row=table_start, column=chart_visits_col, value='Visits 30d')
    ws.cell(row=table_start, column=chart_viewers_col, value='Viewers 30d')
    for c in range(chart_label_col, chart_viewers_col + 1):
        cell = ws.cell(row=table_start, column=c)
        cell.fill = HEADER_FILL
        cell.font = HEADER_FONT
        cell.alignment = Alignment(horizontal='center', vertical='center')
        cell.border = THIN_BORDER

    top_n = min(10, len(rows))
    for i in range(top_n):
        ws.cell(row=table_start + 1 + i, column=chart_label_col, value=rows[i][0])
        ws.cell(row=table_start + 1 + i, column=chart_visits_col, value=rows[i][1])
        ws.cell(row=table_start + 1 + i, column=chart_viewers_col, value=rows[i][2])

    if top_n > 0:
        bar = BarChart()
        bar.type = 'bar'
        bar.style = 11
        bar.title = 'Top 10 Sites — 30-Day Visits & Unique Viewers'
        bar.y_axis.title = 'Site'
        bar.x_axis.title = 'Count (last 30 days)'
        data = Reference(ws, min_col=chart_visits_col, max_col=chart_viewers_col,
                         min_row=table_start, max_row=table_start + top_n)
        cats = Reference(ws, min_col=chart_label_col,
                         min_row=table_start + 1, max_row=table_start + top_n)
        bar.add_data(data, titles_from_data=True)
        bar.set_categories(cats)
        bar.height = 12
        bar.width = 24
        bar.dataLabels = DataLabelList(showVal=True)
        force_axes_visible(bar, x_format='#,##0', y_format='General')
        ws.add_chart(bar, f'{get_column_letter(chart_col)}{last_row + 4}')

    # Lower-traffic granular chart
    lower = [r for r in rows[top_n:] if r[1] > 0 or r[2] > 0]
    if lower:
        low_start = last_row + 32
        ws.cell(row=low_start, column=chart_label_col, value='Site (Rank 11+)')
        ws.cell(row=low_start, column=chart_visits_col, value='Visits 30d')
        ws.cell(row=low_start, column=chart_viewers_col, value='Viewers 30d')
        for c in range(chart_label_col, chart_viewers_col + 1):
            cell = ws.cell(row=low_start, column=c)
            cell.fill = HEADER_FILL
            cell.font = HEADER_FONT
            cell.alignment = Alignment(horizontal='center', vertical='center')
            cell.border = THIN_BORDER

        for i, r in enumerate(lower):
            ws.cell(row=low_start + 1 + i, column=chart_label_col, value=r[0])
            ws.cell(row=low_start + 1 + i, column=chart_visits_col, value=r[1])
            ws.cell(row=low_start + 1 + i, column=chart_viewers_col, value=r[2])
        low_end = low_start + len(lower)

        low_bar = BarChart()
        low_bar.type = 'bar'
        low_bar.style = 12
        low_bar.title = 'Lower-Traffic Sites — Granular 30-Day Visits & Viewers'
        low_bar.y_axis.title = 'Site'
        low_bar.x_axis.title = 'Count (last 30 days)'
        data = Reference(ws, min_col=chart_visits_col, max_col=chart_viewers_col,
                         min_row=low_start, max_row=low_end)
        cats = Reference(ws, min_col=chart_label_col,
                         min_row=low_start + 1, max_row=low_end)
        low_bar.add_data(data, titles_from_data=True)
        low_bar.set_categories(cats)
        low_bar.height = max(10, min(28, 0.6 * len(lower) + 5))
        low_bar.width = 24
        low_bar.dataLabels = DataLabelList(showVal=True)
        force_axes_visible(low_bar, x_format='#,##0', y_format='General')
        ws.add_chart(low_bar, f'{get_column_letter(chart_col)}{last_row + 36}')

    # Callout: Sites With No 30-Day Activity
    inactive_rows = [r for r in rows if r[1] == 0]
    callout_row = last_row + 4

    hdr = ws.cell(row=callout_row, column=1,
                  value=f'⚠ Sites With No 30-Day Activity  ({len(inactive_rows)} of {len(rows)})')
    hdr.font = Font(bold=True, size=13, color='C00000', name='Calibri')
    ws.merge_cells(start_row=callout_row, start_column=1, end_row=callout_row, end_column=8)

    sub = ws.cell(row=callout_row + 1, column=1,
                  value='Review these sites with business owners. Consider archiving, reassigning ownership, or deleting.')
    sub.font = SUBTITLE_FONT
    ws.merge_cells(start_row=callout_row + 1, start_column=1, end_row=callout_row + 1, end_column=8)

    inactive_hdr_row = callout_row + 3
    inactive_headers = ['Site', 'Visits 90d', 'Viewers 90d', 'Visits All-time', 'Active Days (90d)', 'Suggested Action']
    for i, h in enumerate(inactive_headers, 1):
        ws.cell(row=inactive_hdr_row, column=i, value=h)
    style_header_row(ws, inactive_hdr_row, len(inactive_headers))

    site_lookup = {s['site']: s for s in sites_data}
    if inactive_rows:
        for i, r in enumerate(inactive_rows):
            site = r[0]
            s = site_lookup.get(site)
            days_active = 0
            if s is not None and not s['overall']['daily'].empty:
                days_active = int((s['overall']['daily']['site_visits'] > 0).sum())
            viewers_90d = (s['overall']['viewers_90d'] or 0) if s else 0
            row_idx = inactive_hdr_row + 1 + i
            if r[6] == 0 and r[5] == 0:
                action = 'Archive or delete — no lifetime usage'
            elif r[5] == 0:
                action = 'Confirm business owner — no 90d usage'
            else:
                action = 'Validate purpose — cold in last 30d'

            values = [site, r[5], viewers_90d, r[6], days_active, action]
            for c, val in enumerate(values, 1):
                cell = ws.cell(row=row_idx, column=c, value=val)
                cell.border = THIN_BORDER
                if c in (2, 3, 4, 5):
                    cell.number_format = '#,##0'
                    cell.alignment = Alignment(horizontal='right')
                if c == 1:
                    cell.font = Font(bold=True)
                if c == 6:
                    cell.font = Font(italic=True, color='C00000')
                if i % 2 == 1:
                    cell.fill = PatternFill('solid', start_color='FFF4F4')
    else:
        ok = ws.cell(row=inactive_hdr_row + 1, column=1,
                     value='✓ All sites had at least some activity in the last 30 days.')
        ok.font = Font(italic=True, color='548235')
        ws.merge_cells(start_row=inactive_hdr_row + 1, start_column=1,
                       end_row=inactive_hdr_row + 1, end_column=6)


def write_trends_sheet(wb, sites_data):
    """Tenant-wide daily trend (sum of all sites)."""
    ws = wb.create_sheet('Daily Trends')

    ws['A1'] = 'Daily Activity — All Sites Combined'
    ws['A1'].font = TITLE_FONT
    ws.merge_cells('A1:E1')

    all_daily = []
    for s in sites_data:
        if not s['overall']['daily'].empty:
            all_daily.append(s['overall']['daily'])
    if not all_daily:
        ws['A3'] = 'No daily data available.'
        return
    combined = pd.concat(all_daily, ignore_index=True)
    combined = combined.groupby('date', as_index=False).sum()
    combined = combined.sort_values('date')

    headers = ['Date', 'Unique Viewers (sum)', 'Site Visits (sum)']
    for i, h in enumerate(headers, 1):
        ws.cell(row=3, column=i, value=h)
    style_header_row(ws, 3, len(headers))

    for i, row in enumerate(combined.itertuples(index=False), 1):
        ws.cell(row=3 + i, column=1, value=row.date).number_format = 'yyyy-mm-dd'
        ws.cell(row=3 + i, column=2, value=int(row.unique_viewers)).number_format = '#,##0'
        ws.cell(row=3 + i, column=3, value=int(row.site_visits)).number_format = '#,##0'

    last = 3 + len(combined)
    autosize_columns(ws, {'A': 14, 'B': 22, 'C': 18})
    ws.freeze_panes = 'A4'

    line = LineChart()
    line.title = 'Daily Site Visits — Tenant-Wide'
    line.x_axis.title = 'Date'
    line.y_axis.title = 'Visits'
    line.style = 12
    line.height = 10
    line.width = 24
    data = Reference(ws, min_col=3, min_row=3, max_row=last, max_col=3)
    cats = Reference(ws, min_col=1, min_row=4, max_row=last)
    line.add_data(data, titles_from_data=True)
    line.set_categories(cats)
    force_axes_visible(line, x_format='mmm d', y_format='#,##0')
    ws.add_chart(line, 'E3')

    line2 = LineChart()
    line2.title = 'Daily Unique Viewers — Tenant-Wide'
    line2.x_axis.title = 'Date'
    line2.y_axis.title = 'Viewers'
    line2.style = 13
    line2.height = 10
    line2.width = 24
    data2 = Reference(ws, min_col=2, min_row=3, max_row=last, max_col=2)
    line2.add_data(data2, titles_from_data=True)
    line2.set_categories(cats)
    force_axes_visible(line2, x_format='mmm d', y_format='#,##0')
    ws.add_chart(line2, 'E22')


def write_devices_sheet(wb, sites_data):
    ws = wb.create_sheet('Devices')

    ws['A1'] = 'Device Mix — 90-Day Visits by Type'
    ws['A1'].font = TITLE_FONT
    ws.merge_cells('A1:G1')

    headers = ['Site', 'Desktop', 'Mobile App', 'Mobile Web', 'Tablet', 'Other', 'Total']
    for i, h in enumerate(headers, 1):
        ws.cell(row=3, column=i, value=h)
    style_header_row(ws, 3, len(headers))

    for i, s in enumerate(sites_data, 1):
        d = s['devices']
        excel_row = 3 + i
        ws.cell(row=excel_row, column=1, value=s['site'])
        ws.cell(row=excel_row, column=2, value=d['desktop'])
        ws.cell(row=excel_row, column=3, value=d['mobile_app'])
        ws.cell(row=excel_row, column=4, value=d['mobile_web'])
        ws.cell(row=excel_row, column=5, value=d['tablet'])
        ws.cell(row=excel_row, column=6, value=d['other'])
        ws.cell(row=excel_row, column=7, value=f'=SUM(B{excel_row}:F{excel_row})')
        for c in range(2, 8):
            ws.cell(row=excel_row, column=c).number_format = '#,##0'

    last = 3 + len(sites_data)
    total_row = last + 1
    ws.cell(row=total_row, column=1, value='TOTAL').font = Font(bold=True)
    for col_idx in range(2, 8):
        cl = get_column_letter(col_idx)
        cell = ws.cell(row=total_row, column=col_idx,
                       value=f'=SUM({cl}4:{cl}{last})')
        cell.font = Font(bold=True)
        cell.number_format = '#,##0'
        cell.fill = PatternFill('solid', start_color='E7E6E6')

    autosize_columns(ws, {'A': 42, 'B': 12, 'C': 14, 'D': 14, 'E': 12,
                           'F': 12, 'G': 12})
    ws.freeze_panes = 'B4'

    pie = PieChart()
    pie.title = 'Tenant-Wide Device Mix'
    labels = Reference(ws, min_col=2, max_col=6, min_row=3, max_row=3)
    data = Reference(ws, min_col=2, max_col=6, min_row=total_row, max_row=total_row)
    pie.add_data(data, titles_from_data=False)
    pie.set_categories(labels)
    pie.dataLabels = DataLabelList(showPercent=True)
    pie.height = 10
    pie.width = 14
    ws.add_chart(pie, f'I3')


def write_popular_sheet(wb, sites_data):
    ws = wb.create_sheet('Popular Content')
    ws['A1'] = 'Top Content Across All Sites (last 7 days)'
    ws['A1'].font = TITLE_FONT
    ws.merge_cells('A1:E1')

    all_content = []
    for s in sites_data:
        if not s['popular'].empty:
            all_content.append(s['popular'])
    if not all_content:
        ws['A3'] = 'No popular content data available.'
        return
    combined = pd.concat(all_content, ignore_index=True)
    combined = combined.sort_values('visits_7d', ascending=False).head(50)

    headers = ['Site', 'Content', 'Type', 'Viewers (7d)', 'Visits (7d)']
    for i, h in enumerate(headers, 1):
        ws.cell(row=3, column=i, value=h)
    style_header_row(ws, 3, len(headers))

    for i, row in enumerate(combined.itertuples(index=False), 1):
        ws.cell(row=3 + i, column=1, value=row.site)
        ws.cell(row=3 + i, column=2, value=row.content)
        ws.cell(row=3 + i, column=3, value=row.type)
        ws.cell(row=3 + i, column=4, value=int(row.viewers_7d) if pd.notna(row.viewers_7d) else 0).number_format = '#,##0'
        ws.cell(row=3 + i, column=5, value=int(row.visits_7d) if pd.notna(row.visits_7d) else 0).number_format = '#,##0'

    autosize_columns(ws, {'A': 32, 'B': 50, 'C': 14, 'D': 14, 'E': 14})
    ws.freeze_panes = 'A4'


def write_per_site_sheets(wb, sites_data):
    """One sheet per site with its own daily trend chart."""
    used = set()
    for s in sites_data:
        if s['overall']['daily'].empty:
            continue
        base = re.sub(r'[\[\]:*?/\\\']', '_', s['site'])[:31] or 'Site'
        name = base
        n = 2
        while name in used:
            suffix = f'_{n}'
            name = base[:31 - len(suffix)] + suffix
            n += 1
        used.add(name)
        ws = wb.create_sheet(name)

        ws['A1'] = s['site']
        ws['A1'].font = TITLE_FONT
        ws['A2'] = f"7d: {s['overall']['visits_7d'] or 0:,} visits • 30d: {s['overall']['visits_30d'] or 0:,} • 90d: {s['overall']['visits_90d'] or 0:,}"
        ws['A2'].font = SUBTITLE_FONT

        headers = ['Date', 'Unique Viewers', 'Site Visits']
        for i, h in enumerate(headers, 1):
            ws.cell(row=4, column=i, value=h)
        style_header_row(ws, 4, len(headers))

        for i, row in enumerate(s['overall']['daily'].itertuples(index=False), 1):
            ws.cell(row=4 + i, column=1, value=row.date).number_format = 'yyyy-mm-dd'
            ws.cell(row=4 + i, column=2, value=int(row.unique_viewers)).number_format = '#,##0'
            ws.cell(row=4 + i, column=3, value=int(row.site_visits)).number_format = '#,##0'

        last = 4 + len(s['overall']['daily'])
        autosize_columns(ws, {'A': 14, 'B': 16, 'C': 16})
        ws.freeze_panes = 'A5'

        line = LineChart()
        line.title = f'{s["site"]} — Daily Activity'
        line.x_axis.title = 'Date'
        line.y_axis.title = 'Count'
        line.style = 12
        line.height = 10
        line.width = 22
        viewers_ref = Reference(ws, min_col=2, min_row=4, max_row=last, max_col=2)
        visits_ref = Reference(ws, min_col=3, min_row=4, max_row=last, max_col=3)
        cats = Reference(ws, min_col=1, min_row=5, max_row=last)
        line.add_data(viewers_ref, titles_from_data=True)
        line.add_data(visits_ref, titles_from_data=True)
        line.set_categories(cats)
        force_axes_visible(line, x_format='mmm d', y_format='#,##0')
        ws.add_chart(line, 'E4')


# ---------- PowerPoint Output ----------

def _fmt_int(v):
    try:
        return f'{int(v):,}'
    except Exception:
        return '0'


def create_powerpoint_deck(pptx_path, sites_data):
    """Create a management-ready PowerPoint deck focused on 30-day activity."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.chart.data import CategoryChartData

    PRIMARY   = RGBColor(31, 78, 121)
    SECONDARY = RGBColor(91, 155, 213)
    ACCENT    = RGBColor(112, 173, 71)
    CAUTION   = RGBColor(192, 0, 0)
    DARK      = RGBColor(64, 64, 64)
    LIGHT     = RGBColor(242, 246, 250)
    WHITE     = RGBColor(255, 255, 255)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def blank():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        bar = s.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.12))
        bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background()
        return s

    def title(s, t, sub=None):
        b = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.5))
        p = b.text_frame.paragraphs[0]
        p.text = t; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = PRIMARY
        if sub:
            sb = s.shapes.add_textbox(Inches(0.58), Inches(0.85), Inches(12), Inches(0.3))
            sp = sb.text_frame.paragraphs[0]
            sp.text = sub; sp.font.size = Pt(11); sp.font.color.rgb = DARK

    def footer(s):
        f = s.shapes.add_textbox(Inches(0.55), Inches(7.15), Inches(12.2), Inches(0.22))
        p = f.text_frame.paragraphs[0]
        p.text = f'Generated {datetime.now():%B %d, %Y} • SharePoint usage analytics'
        p.font.size = Pt(8); p.font.color.rgb = RGBColor(120, 120, 120)

    def kpi(s, x, y, w, label, value, color=PRIMARY, fill=LIGHT):
        card = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(1.0))
        card.fill.solid(); card.fill.fore_color.rgb = fill
        card.line.color.rgb = RGBColor(210, 220, 230)
        lb = s.shapes.add_textbox(Inches(x+0.1), Inches(y+0.08), Inches(w-0.2), Inches(0.28))
        p = lb.text_frame.paragraphs[0]
        p.text = label.upper(); p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = DARK
        vb = s.shapes.add_textbox(Inches(x+0.1), Inches(y+0.38), Inches(w-0.2), Inches(0.55))
        vp = vb.text_frame.paragraphs[0]
        vp.text = str(value); vp.font.size = Pt(24); vp.font.bold = True; vp.font.color.rgb = color

    # --- Aggregations ---
    ranked = sorted(
        [{
            'site': s['site'],
            'visits_30d': s['overall']['visits_30d'] or 0,
            'viewers_30d': s['overall']['viewers_30d'] or 0,
            'visits_90d': s['overall']['visits_90d'] or 0,
            'visits_all': s['overall']['visits_all'] or 0,
        } for s in sites_data],
        key=lambda r: r['visits_30d'], reverse=True,
    )
    total_visits_30d = sum(r['visits_30d'] for r in ranked)
    total_viewers_30d = sum(r['viewers_30d'] for r in ranked)
    total_visits_7d = sum((s['overall']['visits_7d'] or 0) for s in sites_data)
    total_viewers_7d = sum((s['overall']['viewers_7d'] or 0) for s in sites_data)
    active = sum(1 for r in ranked if r['visits_30d'] > 0)
    inactive = [r for r in ranked if r['visits_30d'] == 0]
    top10 = ranked[:10]
    lower = [r for r in ranked[10:] if r['visits_30d'] > 0 or r['viewers_30d'] > 0]

    # Popular content across all sites
    popular_rows = []
    for site in sites_data:
        if not site['popular'].empty:
            for row in site['popular'].itertuples(index=False):
                try:
                    v7 = int(row.visits_7d) if pd.notna(row.visits_7d) else 0
                    u7 = int(row.viewers_7d) if pd.notna(row.viewers_7d) else 0
                except Exception:
                    v7, u7 = 0, 0
                popular_rows.append({
                    'site': str(row.site),
                    'content': str(row.content),
                    'type': str(row.type),
                    'viewers_7d': u7,
                    'visits_7d': v7,
                })
    popular_rows.sort(key=lambda r: r['visits_7d'], reverse=True)

    # Device totals across all sites
    device_totals = {
        'Desktop':    sum(s['devices']['desktop']    for s in sites_data),
        'Mobile App': sum(s['devices']['mobile_app'] for s in sites_data),
        'Mobile Web': sum(s['devices']['mobile_web'] for s in sites_data),
        'Tablet':     sum(s['devices']['tablet']     for s in sites_data),
        'Other':      sum(s['devices']['other']      for s in sites_data),
    }

    # --- Slide 1: Title ---
    s = blank()
    t = s.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.8), Inches(0.9))
    p = t.text_frame.paragraphs[0]
    p.text = 'SharePoint Site Usage Review'
    p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = PRIMARY
    st = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.3), Inches(0.55))
    sp = st.text_frame.paragraphs[0]
    sp.text = f'Executive summary for {len(sites_data)} SharePoint sites • {datetime.now():%B %d, %Y}'
    sp.font.size = Pt(16); sp.font.color.rgb = DARK
    kpi(s, 0.8,  4.1, 2.9, '30-Day Visits',        _fmt_int(total_visits_30d),  PRIMARY,  LIGHT)
    kpi(s, 3.95, 4.1, 2.9, '30-Day Unique Viewers', _fmt_int(total_viewers_30d), PRIMARY,  LIGHT)
    kpi(s, 7.1,  4.1, 2.9, 'Active Sites (30d)',   f'{active}/{len(sites_data)}', ACCENT,  RGBColor(240,248,240))
    kpi(s, 10.25, 4.1, 2.9, 'Inactive Sites (30d)', _fmt_int(len(inactive)), CAUTION, RGBColor(253,240,240))
    footer(s)

    # --- Slide 2: Executive findings ---
    s = blank(); title(s, 'Executive Findings', 'What management should focus on')
    kpi(s, 0.65, 1.35, 2.4, '30-Day Visits',         _fmt_int(total_visits_30d),  PRIMARY, LIGHT)
    kpi(s, 3.20, 1.35, 2.4, '30-Day Unique Viewers', _fmt_int(total_viewers_30d), PRIMARY, LIGHT)
    kpi(s, 5.75, 1.35, 2.4, '7-Day Visits',          _fmt_int(total_visits_7d),   SECONDARY)
    kpi(s, 8.30, 1.35, 2.4, '7-Day Unique Viewers',  _fmt_int(total_viewers_7d),  SECONDARY)
    kpi(s, 10.85, 1.35, 2.1, 'Inactive Sites', _fmt_int(len(inactive)), CAUTION, RGBColor(253,240,240))
    bullets = []
    if top10:
        bullets.append(f"{top10[0]['site']} leads with {_fmt_int(top10[0]['visits_30d'])} 30-day visits and {_fmt_int(top10[0]['viewers_30d'])} unique viewers.")
    bullets.append(f"{len(inactive)} of {len(ranked)} sites had zero 30-day activity — candidates for review, archival, or cleanup.")
    bullets.append("Usage is concentrated: the top 10 sites drive most 30-day traffic; lower-traffic sites need targeted promotion or retirement.")
    bullets.append("Popular content is heavily document-driven, useful for identifying operationally critical files.")
    box = s.shapes.add_textbox(Inches(0.85), Inches(2.7), Inches(11.8), Inches(3.8))
    tf = box.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'• {b}'; p.font.size = Pt(18); p.font.color.rgb = DARK
    footer(s)

    # --- Slide 3: Top 10 ---
    s = blank(); title(s, 'Top 10 Sites — 30-Day Visits & Unique Viewers', 'Where SharePoint usage is concentrated')
    cd = CategoryChartData()
    cd.categories = [r['site'][:28] for r in top10]
    cd.add_series('30-Day Visits',  [r['visits_30d']  for r in top10])
    cd.add_series('30-Day Viewers', [r['viewers_30d'] for r in top10])
    ch = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                            Inches(0.55), Inches(1.25), Inches(12.2), Inches(5.55), cd).chart
    ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout = False
    ch.value_axis.tick_labels.number_format = '#,##0'
    for series in ch.series:
        series.data_labels.show_value = True
        series.data_labels.number_format = '#,##0'
        series.data_labels.font.size = Pt(9)
    footer(s)

    # --- Slide 4: Lower-traffic sites ---
    if lower:
        s = blank(); title(s, 'Lower-Traffic Sites — Granular 30-Day View',
                           'Smaller-scale sites where every visit and viewer matters')
        cd = CategoryChartData()
        cd.categories = [r['site'][:28] for r in lower]
        cd.add_series('30-Day Visits',  [r['visits_30d']  for r in lower])
        cd.add_series('30-Day Viewers', [r['viewers_30d'] for r in lower])
        ch = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                                Inches(0.55), Inches(1.25), Inches(12.2), Inches(5.55), cd).chart
        ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout = False
        ch.value_axis.tick_labels.number_format = '#,##0'
        for series in ch.series:
            series.data_labels.show_value = True
            series.data_labels.number_format = '#,##0'
            series.data_labels.font.size = Pt(9)
        footer(s)

    # --- Slide 5: Popular Content across sites (7-day) ---
    if popular_rows:
        s = blank()
        title(s, 'Most-Used Content Across Sites',
              'Top content by 7-day visits — which files staff actually use')
        top_content = popular_rows[:12]
        cols = 5
        rows = len(top_content) + 1
        table = s.shapes.add_table(rows, cols, Inches(0.4), Inches(1.25),
                                   Inches(12.55), Inches(5.6)).table
        table.columns[0].width = Inches(2.6)   # Site
        table.columns[1].width = Inches(5.35)  # Content
        table.columns[2].width = Inches(1.4)   # Type
        table.columns[3].width = Inches(1.4)   # Viewers 7d
        table.columns[4].width = Inches(1.8)   # Visits 7d
        hdrs = ['Site', 'Content', 'Type', 'Viewers 7d', 'Visits 7d']
        for c, h in enumerate(hdrs):
            cell = table.cell(0, c); cell.text = h
            cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = WHITE; para.font.bold = True; para.font.size = Pt(11)
        for r_idx, row in enumerate(top_content, 1):
            vals = [
                row['site'][:32],
                row['content'][:68],
                row['type'][:12],
                _fmt_int(row['viewers_7d']),
                _fmt_int(row['visits_7d']),
            ]
            for c, v in enumerate(vals):
                cell = table.cell(r_idx, c); cell.text = v
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                if r_idx % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
                if c == 4:  # Visits column emphasis
                    para.font.bold = True
                    para.font.color.rgb = PRIMARY
        footer(s)

    # --- Slide 6: Device Mix (90-day) ---
    active_devices = {k: v for k, v in device_totals.items() if v > 0}
    if active_devices:
        s = blank()
        title(s, 'Device Mix — 90-Day Visits by Type',
              'How staff are accessing SharePoint')
        cd = CategoryChartData()
        cd.categories = list(active_devices.keys())
        cd.add_series('Visits', list(active_devices.values()))
        ch = s.shapes.add_chart(XL_CHART_TYPE.PIE,
                                Inches(0.5), Inches(1.35),
                                Inches(6.6), Inches(5.4), cd).chart
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        plot = ch.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.show_percentage = True
        dl.show_category_name = False
        dl.font.size = Pt(12)
        dl.font.bold = True

        total_devices = sum(active_devices.values())
        rows = len(active_devices) + 2  # header + rows + total
        table = s.shapes.add_table(rows, 3, Inches(7.5), Inches(1.6),
                                   Inches(5.4), Inches(4.5)).table
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(1.7)
        table.columns[2].width = Inches(1.7)
        hdrs = ['Device', 'Visits (90d)', 'Share']
        for c, h in enumerate(hdrs):
            cell = table.cell(0, c); cell.text = h
            cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = WHITE; para.font.bold = True; para.font.size = Pt(12)
        for r_idx, (dev, count) in enumerate(active_devices.items(), 1):
            share = (count / total_devices * 100) if total_devices else 0
            vals = [dev, _fmt_int(count), f'{share:.1f}%']
            for c, v in enumerate(vals):
                cell = table.cell(r_idx, c); cell.text = v
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(12)
                if r_idx % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
        trow = len(active_devices) + 1
        totals_row = ['TOTAL', _fmt_int(total_devices), '100.0%']
        for c, v in enumerate(totals_row):
            cell = table.cell(trow, c); cell.text = v
            cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = WHITE; para.font.bold = True; para.font.size = Pt(12)
        footer(s)

    # --- Slide 7: Sites with NO 30-day activity ---
    s = blank(); title(s, 'Sites With No 30-Day Activity',
                       f'{len(inactive)} of {len(ranked)} sites had zero visits in the last 30 days')
    if inactive:
        cols = 4
        rows = min(len(inactive), 14) + 1
        table = s.shapes.add_table(rows, cols, Inches(0.55), Inches(1.25),
                                   Inches(12.2), Inches(5.6)).table
        hdrs = ['Site', 'Visits 90d', 'Visits All-time', 'Suggested Action']
        for c, h in enumerate(hdrs):
            cell = table.cell(0, c); cell.text = h
            cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(11)
        for r_idx, r in enumerate(inactive[:14], 1):
            if r['visits_all'] == 0 and r['visits_90d'] == 0:
                action = 'Archive or delete — no lifetime usage'
            elif r['visits_90d'] == 0:
                action = 'Confirm owner — no 90d usage'
            else:
                action = 'Validate purpose — cold in last 30d'
            vals = [r['site'][:38], _fmt_int(r['visits_90d']),
                    _fmt_int(r['visits_all']), action]
            for c, v in enumerate(vals):
                cell = table.cell(r_idx, c); cell.text = v
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                if c == 3:
                    cell.text_frame.paragraphs[0].font.color.rgb = CAUTION
                    cell.text_frame.paragraphs[0].font.italic = True
    else:
        b = s.shapes.add_textbox(Inches(0.85), Inches(3), Inches(11.5), Inches(1))
        p = b.text_frame.paragraphs[0]
        p.text = '✓ All sites had at least some 30-day activity.'
        p.font.size = Pt(24); p.font.color.rgb = ACCENT; p.font.bold = True
    footer(s)

    prs.save(pptx_path)


# ---------- Main ----------

def main():
    parser = argparse.ArgumentParser(
        description='Aggregate SharePoint Site Analytics exports into an executive Excel dashboard (and optional PowerPoint deck).'
    )
    parser.add_argument('input_folder', help='Folder containing SharePoint analytics .xlsx exports')
    parser.add_argument('output_file', nargs='?', help='Output Excel workbook path')
    parser.add_argument('--pptx', dest='pptx_file',
                        help='Optional PowerPoint deck output path (requires python-pptx)')
    parser.add_argument('--skip-site-sheets', action='store_true',
                        help='Skip per-site worksheets to keep the workbook smaller')
    args = parser.parse_args()

    in_dir = Path(args.input_folder)
    out_path = Path(args.output_file) if args.output_file else \
        in_dir / f'SharePoint-Usage-Dashboard-{datetime.now():%Y-%m-%d}.xlsx'

    if not in_dir.is_dir():
        print(f'Not a directory: {in_dir}')
        sys.exit(1)

    files = sorted([p for p in in_dir.glob('*.xlsx')
                    if not p.name.startswith('~')
                    and 'Dashboard' not in p.name
                    and 'Aggregate' not in p.name])
    print(f'Found {len(files)} files in {in_dir}')

    sites = []
    for p in files:
        print(f'  reading {p.name}')
        data = load_site(p)
        if data:
            sites.append(data)

    if not sites:
        print('No usable data found.')
        sys.exit(1)

    print(f'\nBuilding dashboard with {len(sites)} sites...')
    wb = Workbook()
    wb.remove(wb.active)
    write_summary_sheet(wb, sites)
    write_trends_sheet(wb, sites)
    write_devices_sheet(wb, sites)
    write_popular_sheet(wb, sites)
    if not args.skip_site_sheets:
        write_per_site_sheets(wb, sites)

    wb.save(out_path)
    print(f'\nSaved Excel dashboard: {out_path}')

    if args.pptx_file:
        pptx_path = Path(args.pptx_file)
        try:
            create_powerpoint_deck(pptx_path, sites)
            print(f'Saved PowerPoint deck:  {pptx_path}')
        except ImportError:
            print('! PowerPoint output requires python-pptx. Install with: pip install python-pptx')


if __name__ == '__main__':
    main()